"""Managed upload registry and local processing pipeline for source documents."""

from __future__ import annotations

import csv
import datetime as dt
import hashlib
import json
import os
import re
import shutil
import sqlite3
import subprocess
import tempfile
import textwrap
import time
import zipfile
from contextlib import contextmanager
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from lxml import etree
from openpyxl import load_workbook
from pptx import Presentation

try:
    from pypdf import PdfReader
except ImportError:  # pragma: no cover - dependency can be absent in minimal runtimes
    PdfReader = None

from ops.docs import service
from ops.shared.io import sha256_file, write_json_atomic

UPLOADS_DIR = service.DOCS_DIR / "uploads"
UPLOADS_DIR.mkdir(parents=True, exist_ok=True)
DERIVED_DIR = UPLOADS_DIR / ".derived"
DERIVED_DIR.mkdir(parents=True, exist_ok=True)

UPLOAD_ID_PATTERN = re.compile(r"^UPL-[a-f0-9]{32}$")
SESSION_ID_PATTERN = re.compile(r"^[A-Za-z0-9_-]{8,128}$")
SHA256_PATTERN = re.compile(r"^[a-f0-9]{64}$")
DOCUMENT_ID_PATTERN = re.compile(r"^DOC-[a-f0-9]{32}$")

SCHEMA_VERSION = "3"
CURRENT_EXTRACTOR_VERSION = "local-v1"
MAX_TEXT_CHARS = 120_000
MAX_MARKDOWN_CHARS = 120_000
MAX_WARNINGS = 32
MAX_ERROR_CHARS = 800
DEFAULT_JOB_MAX_ATTEMPTS = 5
DEFAULT_JOB_RETRY_SECONDS = 60
DEFAULT_LEASE_SECONDS = 120
MAX_LINKED_UPLOADS_PER_SESSION = 100
MAX_LINKED_UPLOADS_IN_SESSION_DETAIL = 8

TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".markdown",
    ".yaml",
    ".yml",
}
DELIMITED_EXTENSIONS = {
    ".csv": ",",
    ".tsv": "\t",
}
MODERN_ARCHIVE_EXTENSIONS = {
    ".docx",
    ".xlsx",
    ".pptx",
}
LIBREOFFICE_EXTENSIONS = {
    ".doc",
    ".xls",
    ".ppt",
    ".odt",
    ".ods",
    ".odp",
}

_WORD_RE = re.compile(r"[A-Za-z][A-Za-z0-9'-]{1,31}")
_SENTENCE_SPLIT_RE = re.compile(r"(?<=[.!?])\s+")
_RTF_CONTROL_RE = re.compile(r"\\[a-zA-Z]+\d* ?")
_RTF_HEX_RE = re.compile(r"\\'[0-9a-fA-F]{2}")
_REPEATED_WS_RE = re.compile(r"[ \t]{2,}")


@dataclass(frozen=True)
class ExtractionResult:
    status: str
    markdown: str
    text: str
    metadata: dict[str, Any]
    warnings: list[str]
    provenance: dict[str, Any]
    tags: list[str]
    summary: str
    error: str | None = None


@contextmanager
def _db():
    conn = sqlite3.connect(service._DB_PATH)
    conn.row_factory = sqlite3.Row
    try:
        conn.execute(
            """CREATE TABLE IF NOT EXISTS docops_uploads (
                upload_id TEXT NOT NULL,
                session_id TEXT NOT NULL,
                relative_path TEXT NOT NULL,
                name TEXT NOT NULL,
                path TEXT NOT NULL,
                mime_type TEXT NOT NULL,
                byte_size INTEGER NOT NULL,
                sha256 TEXT NOT NULL,
                created_at TEXT DEFAULT CURRENT_TIMESTAMP,
                PRIMARY KEY (upload_id, relative_path)
            )"""
        )
        conn.execute(
            "CREATE INDEX IF NOT EXISTS idx_docops_uploads_session "
            "ON docops_uploads(session_id, created_at)"
        )
        _ensure_upload_schema(conn)
        yield conn
        conn.commit()
    finally:
        conn.close()


def _ensure_upload_schema(conn: sqlite3.Connection) -> None:
    conn.execute(
        """CREATE TABLE IF NOT EXISTS docops_upload_meta (
            key TEXT PRIMARY KEY,
            value TEXT NOT NULL
        )"""
    )
    conn.execute(
        """CREATE TABLE IF NOT EXISTS docops_upload_documents (
            document_id TEXT PRIMARY KEY,
            sha256 TEXT NOT NULL UNIQUE,
            byte_size INTEGER NOT NULL,
            canonical_path TEXT NOT NULL,
            canonical_name TEXT NOT NULL,
            canonical_mime_type TEXT NOT NULL,
            created_at TEXT DEFAULT CURRENT_TIMESTAMP,
            updated_at TEXT DEFAULT CURRENT_TIMESTAMP
        )"""
    )
    conn.execute(
        """CREATE TABLE IF NOT EXISTS docops_upload_instances (
            instance_id TEXT PRIMARY KEY,
            document_id TEXT NOT NULL,
            upload_id TEXT NOT NULL,
            session_id TEXT NOT NULL,
            relative_path TEXT NOT NULL,
            name TEXT NOT NULL,
            path TEXT NOT NULL,
            mime_type TEXT NOT NULL,
            byte_size INTEGER NOT NULL,
            sha256 TEXT NOT NULL,
            created_at TEXT DEFAULT CURRENT_TIMESTAMP,
            UNIQUE (upload_id, relative_path),
            FOREIGN KEY (document_id) REFERENCES docops_upload_documents(document_id)
        )"""
    )
    conn.execute(
        "CREATE INDEX IF NOT EXISTS idx_upload_instances_session "
        "ON docops_upload_instances(session_id, created_at)"
    )
    conn.execute(
        "CREATE INDEX IF NOT EXISTS idx_upload_instances_document "
        "ON docops_upload_instances(document_id, created_at)"
    )
    conn.execute(
        """CREATE TABLE IF NOT EXISTS docops_upload_extractions (
            document_id TEXT NOT NULL,
            extractor_version TEXT NOT NULL,
            status TEXT NOT NULL,
            markdown_text TEXT NOT NULL,
            normalized_text TEXT NOT NULL,
            structural_metadata_json TEXT NOT NULL,
            tags_json TEXT NOT NULL,
            summary_text TEXT NOT NULL,
            warnings_json TEXT NOT NULL,
            provenance_json TEXT NOT NULL,
            derived_root TEXT NOT NULL,
            error_detail TEXT NOT NULL DEFAULT '',
            created_at TEXT DEFAULT CURRENT_TIMESTAMP,
            updated_at TEXT DEFAULT CURRENT_TIMESTAMP,
            PRIMARY KEY (document_id, extractor_version),
            FOREIGN KEY (document_id) REFERENCES docops_upload_documents(document_id)
        )"""
    )
    conn.execute(
        "CREATE INDEX IF NOT EXISTS idx_upload_extractions_status "
        "ON docops_upload_extractions(status, updated_at)"
    )
    conn.execute(
        """CREATE TABLE IF NOT EXISTS docops_upload_jobs (
            job_id TEXT PRIMARY KEY,
            document_id TEXT NOT NULL,
            document_sha256 TEXT NOT NULL,
            extractor_version TEXT NOT NULL,
            status TEXT NOT NULL,
            attempts INTEGER NOT NULL DEFAULT 0,
            max_attempts INTEGER NOT NULL DEFAULT 5,
            lease_owner TEXT,
            lease_expires_at TEXT,
            heartbeat_at TEXT,
            next_attempt_at TEXT NOT NULL,
            last_error TEXT NOT NULL DEFAULT '',
            created_at TEXT DEFAULT CURRENT_TIMESTAMP,
            updated_at TEXT DEFAULT CURRENT_TIMESTAMP,
            completed_at TEXT,
            UNIQUE (document_sha256, extractor_version),
            FOREIGN KEY (document_id) REFERENCES docops_upload_documents(document_id)
        )"""
    )
    conn.execute(
        "CREATE INDEX IF NOT EXISTS idx_upload_jobs_ready "
        "ON docops_upload_jobs(status, next_attempt_at, lease_expires_at)"
    )
    conn.execute(
        """CREATE TABLE IF NOT EXISTS docops_chat_session_upload_documents (
            session_id TEXT NOT NULL,
            document_id TEXT NOT NULL,
            source_session_id TEXT NOT NULL,
            linked_at TEXT NOT NULL,
            PRIMARY KEY (session_id, document_id),
            FOREIGN KEY (document_id) REFERENCES docops_upload_documents(document_id)
        )"""
    )
    conn.execute(
        "CREATE INDEX IF NOT EXISTS idx_chat_session_upload_documents_session "
        "ON docops_chat_session_upload_documents(session_id, linked_at)"
    )
    conn.execute(
        """CREATE TABLE IF NOT EXISTS docops_chat_session_upload_instances (
            session_id TEXT NOT NULL,
            instance_id TEXT NOT NULL,
            document_id TEXT NOT NULL,
            upload_id TEXT NOT NULL,
            relative_path TEXT NOT NULL,
            source_session_id TEXT NOT NULL,
            linked_at TEXT NOT NULL,
            PRIMARY KEY (session_id, instance_id),
            FOREIGN KEY (instance_id) REFERENCES docops_upload_instances(instance_id),
            FOREIGN KEY (document_id) REFERENCES docops_upload_documents(document_id)
        )"""
    )
    conn.execute(
        "CREATE INDEX IF NOT EXISTS idx_chat_session_upload_instances_session "
        "ON docops_chat_session_upload_instances(session_id, linked_at)"
    )

    current = conn.execute(
        "SELECT value FROM docops_upload_meta WHERE key='schema_version'"
    ).fetchone()
    if current and current["value"] == SCHEMA_VERSION:
        return

    legacy_rows = conn.execute(
        "SELECT upload_id, session_id, relative_path, name, path, mime_type, "
        "byte_size, sha256, created_at FROM docops_uploads ORDER BY created_at, rowid"
    ).fetchall()
    for row in legacy_rows:
        document_id = _document_id_from_sha256(row["sha256"])
        conn.execute(
            "INSERT OR IGNORE INTO docops_upload_documents "
            "(document_id, sha256, byte_size, canonical_path, canonical_name, canonical_mime_type, "
            "created_at, updated_at) VALUES (?,?,?,?,?,?,?,?)",
            (
                document_id,
                row["sha256"],
                int(row["byte_size"]),
                row["path"],
                row["name"],
                row["mime_type"],
                row["created_at"],
                row["created_at"],
            ),
        )
        conn.execute(
            "INSERT OR IGNORE INTO docops_upload_instances "
            "(instance_id, document_id, upload_id, session_id, relative_path, name, path, mime_type, "
            "byte_size, sha256, created_at) VALUES (?,?,?,?,?,?,?,?,?,?,?)",
            (
                _instance_id(row["upload_id"], row["relative_path"]),
                document_id,
                row["upload_id"],
                row["session_id"],
                row["relative_path"],
                row["name"],
                row["path"],
                row["mime_type"],
                int(row["byte_size"]),
                row["sha256"],
                row["created_at"],
            ),
        )
        _enqueue_processing_job_conn(
            conn,
            document_id=document_id,
            document_sha256=row["sha256"],
            extractor_version=CURRENT_EXTRACTOR_VERSION,
        )
    conn.execute(
        "INSERT INTO docops_upload_meta (key, value) VALUES ('schema_version', ?) "
        "ON CONFLICT(key) DO UPDATE SET value=excluded.value",
        (SCHEMA_VERSION,),
    )


def _utc_now_sql() -> str:
    return dt.datetime.now(dt.timezone.utc).strftime("%Y-%m-%d %H:%M:%S")


def _plus_seconds_sql(base: str, seconds: int) -> str:
    parsed = dt.datetime.strptime(base, "%Y-%m-%d %H:%M:%S").replace(tzinfo=dt.timezone.utc)
    return (parsed + dt.timedelta(seconds=seconds)).strftime("%Y-%m-%d %H:%M:%S")


def _document_id_from_sha256(sha256_hex: str) -> str:
    if not SHA256_PATTERN.fullmatch(sha256_hex or ""):
        raise ValueError("invalid document sha256")
    return f"DOC-{sha256_hex[:32]}"


def _instance_id(upload_id: str, relative_path: str) -> str:
    digest = hashlib.sha256(f"{upload_id}:{relative_path}".encode("utf-8")).hexdigest()
    return f"UPI-{digest[:32]}"


def _job_id(document_sha256: str, extractor_version: str) -> str:
    digest = hashlib.sha256(f"{document_sha256}:{extractor_version}".encode("utf-8")).hexdigest()
    return f"UPJ-{digest[:32]}"


def _bounded(value: str, *, limit: int) -> str:
    normalized = " ".join(str(value or "").split())
    if len(normalized) <= limit:
        return normalized
    return normalized[:limit]


def _public_upload_from_row(row: sqlite3.Row) -> dict[str, Any]:
    return {
        "upload_id": row["upload_id"],
        "session_id": row["session_id"],
        "saved_path": row["relative_path"],
        "name": row["name"],
        "mime": row["mime_type"],
        "size": row["byte_size"],
        "sha256": row["sha256"],
        "created_at": row["created_at"],
        "document_id": row["document_id"],
        "reused": bool(row["reused"]),
    }


def _validate_registration_inputs(upload_id: str, session_id: str, files: list[dict]) -> None:
    if not UPLOAD_ID_PATTERN.fullmatch(upload_id or ""):
        raise ValueError("invalid upload_id")
    if not SESSION_ID_PATTERN.fullmatch(session_id or ""):
        raise ValueError("invalid session_id")
    if not files:
        raise ValueError("at least one uploaded file is required")


def _chat_session_exists_conn(conn: sqlite3.Connection, session_id: str) -> bool:
    table = conn.execute(
        "SELECT 1 FROM sqlite_master WHERE type='table' AND name='chat_sessions'"
    ).fetchone()
    if table is None:
        return False
    return (
        conn.execute("SELECT 1 FROM chat_sessions WHERE id=?", (session_id,)).fetchone() is not None
    )


def _link_upload_instances_to_chat_session_conn(
    conn: sqlite3.Connection,
    *,
    target_session_id: str,
    source_session_id: str,
    rows: list[dict[str, Any]],
    linked_at: str,
) -> int:
    linked_instances = 0
    for item in rows:
        conn.execute(
            "INSERT OR IGNORE INTO docops_chat_session_upload_documents "
            "(session_id, document_id, source_session_id, linked_at) VALUES (?,?,?,?)",
            (
                target_session_id,
                item["document_id"],
                source_session_id,
                linked_at,
            ),
        )
        cur = conn.execute(
            "INSERT OR IGNORE INTO docops_chat_session_upload_instances "
            "(session_id, instance_id, document_id, upload_id, relative_path, source_session_id, linked_at) "
            "VALUES (?,?,?,?,?,?,?)",
            (
                target_session_id,
                _instance_id(item["upload_id"], item["relative_path"]),
                item["document_id"],
                item["upload_id"],
                item["relative_path"],
                source_session_id,
                linked_at,
            ),
        )
        if cur.rowcount:
            linked_instances += 1
    return linked_instances


def register_uploaded_documents(upload_id: str, session_id: str, files: list[dict]) -> dict:
    """Register persisted upload files, dedupe by SHA-256, and enqueue processing."""
    _validate_registration_inputs(upload_id, session_id, files)

    upload_root = UPLOADS_DIR.resolve()
    rows: list[dict[str, Any]] = []
    seen = set()
    for item in files:
        relative_path = str(item.get("saved_path") or "")
        path = Path(item.get("path") or "")
        if not relative_path or relative_path in seen:
            raise ValueError("uploaded paths must be present and unique")
        seen.add(relative_path)
        try:
            resolved = path.resolve(strict=True)
            resolved.relative_to(upload_root)
        except (OSError, ValueError) as exc:
            raise ValueError("uploaded file is outside the managed upload directory") from exc
        if path.is_symlink() or not resolved.is_file():
            raise ValueError("uploaded file must be a regular file")
        byte_size = resolved.stat().st_size
        sha256 = sha256_file(resolved)
        if byte_size != item.get("size") or sha256 != item.get("sha256"):
            raise ValueError("uploaded file integrity verification failed")
        rows.append(
            {
                "upload_id": upload_id,
                "session_id": session_id,
                "relative_path": relative_path,
                "name": str(item.get("name") or resolved.name),
                "path": str(resolved),
                "mime_type": str(item.get("mime") or "application/octet-stream"),
                "byte_size": byte_size,
                "sha256": sha256,
            }
        )

    now = _utc_now_sql()
    registered_rows: list[dict[str, Any]] = []
    enqueue_pairs: dict[str, str] = {}
    with _db() as conn:
        should_link_chat_session = _chat_session_exists_conn(conn, session_id)
        for row in rows:
            existing_doc = conn.execute(
                "SELECT document_id FROM docops_upload_documents WHERE sha256=?",
                (row["sha256"],),
            ).fetchone()
            reused = existing_doc is not None
            document_id = (
                existing_doc["document_id"]
                if existing_doc
                else _document_id_from_sha256(row["sha256"])
            )
            conn.execute(
                "INSERT OR IGNORE INTO docops_upload_documents "
                "(document_id, sha256, byte_size, canonical_path, canonical_name, canonical_mime_type, "
                "created_at, updated_at) VALUES (?,?,?,?,?,?,?,?)",
                (
                    document_id,
                    row["sha256"],
                    row["byte_size"],
                    row["path"],
                    row["name"],
                    row["mime_type"],
                    now,
                    now,
                ),
            )
            conn.execute(
                "INSERT INTO docops_uploads "
                "(upload_id, session_id, relative_path, name, path, mime_type, byte_size, sha256, created_at) "
                "VALUES (?,?,?,?,?,?,?,?,?)",
                (
                    row["upload_id"],
                    row["session_id"],
                    row["relative_path"],
                    row["name"],
                    row["path"],
                    row["mime_type"],
                    row["byte_size"],
                    row["sha256"],
                    now,
                ),
            )
            conn.execute(
                "INSERT INTO docops_upload_instances "
                "(instance_id, document_id, upload_id, session_id, relative_path, name, path, mime_type, "
                "byte_size, sha256, created_at) VALUES (?,?,?,?,?,?,?,?,?,?,?)",
                (
                    _instance_id(row["upload_id"], row["relative_path"]),
                    document_id,
                    row["upload_id"],
                    row["session_id"],
                    row["relative_path"],
                    row["name"],
                    row["path"],
                    row["mime_type"],
                    row["byte_size"],
                    row["sha256"],
                    now,
                ),
            )
            enqueue_pairs[row["sha256"]] = document_id
            registered_rows.append(
                {
                    **row,
                    "created_at": now,
                    "document_id": document_id,
                    "reused": reused,
                }
            )
        for sha256_hex, document_id in enqueue_pairs.items():
            _enqueue_processing_job_conn(
                conn,
                document_id=document_id,
                document_sha256=sha256_hex,
                extractor_version=CURRENT_EXTRACTOR_VERSION,
            )
        if should_link_chat_session:
            _link_upload_instances_to_chat_session_conn(
                conn,
                target_session_id=session_id,
                source_session_id=session_id,
                rows=registered_rows,
                linked_at=now,
            )
        job_rows = conn.execute(
            "SELECT document_id, status, attempts, max_attempts, updated_at, completed_at "
            "FROM docops_upload_jobs WHERE extractor_version=?",
            (CURRENT_EXTRACTOR_VERSION,),
        ).fetchall()
        jobs_by_document = {row["document_id"]: dict(row) for row in job_rows}

    registered_rows.sort(key=lambda item: item["relative_path"])
    return {
        "upload_id": upload_id,
        "count": len(registered_rows),
        "documents": [
            {
                "upload_id": item["upload_id"],
                "session_id": item["session_id"],
                "saved_path": item["relative_path"],
                "name": item["name"],
                "mime": item["mime_type"],
                "size": item["byte_size"],
                "sha256": item["sha256"],
                "created_at": item["created_at"],
                "document_id": item["document_id"],
                "canonical_document_id": item["document_id"],
                "reused": item["reused"],
                "queue_state": (
                    jobs_by_document.get(item["document_id"], {}).get("status") or "queued"
                ),
                "processing": {
                    "job_status": jobs_by_document.get(item["document_id"], {}).get("status")
                    or "queued",
                    "attempts": jobs_by_document.get(item["document_id"], {}).get("attempts") or 0,
                    "max_attempts": jobs_by_document.get(item["document_id"], {}).get(
                        "max_attempts"
                    )
                    or DEFAULT_JOB_MAX_ATTEMPTS,
                    "updated_at": jobs_by_document.get(item["document_id"], {}).get("updated_at"),
                    "completed_at": jobs_by_document.get(item["document_id"], {}).get(
                        "completed_at"
                    ),
                },
            }
            for item in registered_rows
        ],
    }


def list_uploaded_documents(session_id: str = "", query: str = "", limit: int = 50) -> dict:
    """List uploaded source documents registered with DocOps."""
    if not isinstance(limit, int) or not 1 <= limit <= 100:
        return {"error": "limit must be an integer from 1 to 100"}
    like = f"%{query}%"
    with _db() as conn:
        rows = conn.execute(
            "SELECT i.upload_id, i.session_id, i.relative_path, i.name, i.path, i.mime_type, "
            "i.byte_size, i.sha256, i.created_at, i.document_id, "
            "(SELECT COUNT(*) FROM docops_upload_instances all_i WHERE all_i.document_id=i.document_id) > 1 "
            "AS reused "
            "FROM docops_upload_instances i "
            "WHERE (? = '' OR i.session_id = ?) "
            "AND (i.name LIKE ? OR i.relative_path LIKE ? OR i.upload_id LIKE ?) "
            "ORDER BY i.created_at DESC, i.relative_path LIMIT ?",
            (session_id, session_id, like, like, like, limit),
        ).fetchall()
    return {"count": len(rows), "documents": [_public_upload_from_row(row) for row in rows]}


def _load_single_upload(upload_id: str, saved_path: str, session_id: str = "") -> list[sqlite3.Row]:
    if session_id:
        _validate_session_scope(session_id)
    with _db() as conn:
        if session_id:
            return conn.execute(
                "SELECT i.upload_id, i.session_id, i.relative_path, i.name, i.path, i.mime_type, "
                "i.byte_size, i.sha256, i.created_at, i.document_id, "
                "(SELECT COUNT(*) FROM docops_upload_instances all_i WHERE all_i.document_id=i.document_id) > 1 "
                "AS reused "
                "FROM docops_chat_session_upload_instances l "
                "JOIN docops_upload_instances i ON i.instance_id=l.instance_id "
                "WHERE l.session_id=? AND i.upload_id=? AND (? = '' OR i.relative_path=?) "
                "ORDER BY i.relative_path",
                (session_id, upload_id, saved_path, saved_path),
            ).fetchall()
        return conn.execute(
            "SELECT i.upload_id, i.session_id, i.relative_path, i.name, i.path, i.mime_type, "
            "i.byte_size, i.sha256, i.created_at, i.document_id, "
            "(SELECT COUNT(*) FROM docops_upload_instances all_i WHERE all_i.document_id=i.document_id) > 1 "
            "AS reused "
            "FROM docops_upload_instances i "
            "WHERE i.upload_id=? AND (? = '' OR i.relative_path=?) ORDER BY i.relative_path",
            (upload_id, saved_path, saved_path),
        ).fetchall()


def get_uploaded_document(upload_id: str, saved_path: str = "", session_id: str = "") -> dict:
    """Resolve an uploaded source document and return metadata with bounded preview."""
    rows = _load_single_upload(upload_id, saved_path, session_id=session_id)
    if not rows:
        return {"error": f"unknown uploaded document '{upload_id}'"}
    if len(rows) > 1:
        return {
            "upload_id": upload_id,
            "count": len(rows),
            "documents": [_public_upload_from_row(row) for row in rows],
        }

    row = rows[0]
    metadata = _public_upload_from_row(row)
    path = Path(row["path"])
    try:
        resolved = path.resolve(strict=True)
        resolved.relative_to(UPLOADS_DIR.resolve())
    except (OSError, ValueError):
        return {"error": "uploaded document is missing or outside managed storage"}
    if resolved.stat().st_size != row["byte_size"] or sha256_file(resolved) != row["sha256"]:
        return {"error": "uploaded document failed integrity verification"}

    if resolved.suffix.lower() in {
        ".txt",
        ".md",
        ".markdown",
        ".csv",
        ".tsv",
        ".json",
        ".yaml",
        ".yml",
        ".xml",
    }:
        content = resolved.read_text(encoding="utf-8", errors="replace")
        metadata["content"] = content[:20_000]
        metadata["content_truncated"] = len(content) > 20_000
    else:
        metadata["content"] = None
        metadata["content_truncated"] = False

    with _db() as conn:
        extraction = conn.execute(
            "SELECT extractor_version, status, tags_json, summary_text, warnings_json, provenance_json "
            "FROM docops_upload_extractions WHERE document_id=? "
            "ORDER BY updated_at DESC LIMIT 1",
            (row["document_id"],),
        ).fetchone()
    if extraction:
        metadata["processing"] = {
            "extractor_version": extraction["extractor_version"],
            "status": extraction["status"],
            "tags": json.loads(extraction["tags_json"]),
            "summary": extraction["summary_text"],
            "warnings": json.loads(extraction["warnings_json"]),
            "provenance": json.loads(extraction["provenance_json"]),
        }
    return metadata


def _validate_session_scope(session_id: str) -> None:
    if not SESSION_ID_PATTERN.fullmatch(session_id or ""):
        raise ValueError("invalid session_id")


def _document_processing_snapshot_conn(
    conn: sqlite3.Connection, document_id: str
) -> dict[str, Any]:
    extraction = conn.execute(
        "SELECT extractor_version, status, tags_json, summary_text, warnings_json, provenance_json, "
        "updated_at, created_at FROM docops_upload_extractions WHERE document_id=? "
        "ORDER BY updated_at DESC LIMIT 1",
        (document_id,),
    ).fetchone()
    job = conn.execute(
        "SELECT job_id, status, attempts, max_attempts, last_error, next_attempt_at, "
        "updated_at, created_at, completed_at "
        "FROM docops_upload_jobs WHERE document_id=? AND extractor_version=? LIMIT 1",
        (document_id, CURRENT_EXTRACTOR_VERSION),
    ).fetchone()
    return {
        "queue_state": (job["status"] if job else "queued"),
        "job": (
            {
                "job_id": job["job_id"],
                "status": job["status"],
                "attempts": int(job["attempts"]),
                "max_attempts": int(job["max_attempts"]),
                "last_error": job["last_error"],
                "next_attempt_at": job["next_attempt_at"],
                "updated_at": job["updated_at"],
                "created_at": job["created_at"],
                "completed_at": job["completed_at"],
            }
            if job
            else None
        ),
        "extraction": (
            {
                "extractor_version": extraction["extractor_version"],
                "status": extraction["status"],
                "tags": json.loads(extraction["tags_json"] or "[]"),
                "summary": extraction["summary_text"] or "",
                "warnings": json.loads(extraction["warnings_json"] or "[]"),
                "provenance": json.loads(extraction["provenance_json"] or "{}"),
                "updated_at": extraction["updated_at"],
                "created_at": extraction["created_at"],
            }
            if extraction
            else None
        ),
    }


def _session_upload_rows(
    conn: sqlite3.Connection, session_id: str, *, query: str = "", limit: int = 50
) -> list[sqlite3.Row]:
    like = f"%{query}%"
    return conn.execute(
        "SELECT i.instance_id, i.document_id, i.upload_id, i.session_id AS upload_session_id, "
        "i.relative_path, i.name, i.path, i.mime_type, i.byte_size, i.sha256, i.created_at AS upload_created_at, "
        "l.linked_at, l.source_session_id, "
        "(SELECT COUNT(*) FROM docops_upload_instances all_i WHERE all_i.document_id = i.document_id) > 1 AS reused "
        "FROM docops_chat_session_upload_instances l "
        "JOIN docops_upload_instances i ON i.instance_id = l.instance_id "
        "WHERE l.session_id=? "
        "AND (? = '' OR i.name LIKE ? OR i.relative_path LIKE ? OR i.upload_id LIKE ? OR i.document_id LIKE ?) "
        "ORDER BY l.linked_at DESC, i.relative_path ASC LIMIT ?",
        (session_id, query, like, like, like, like, limit),
    ).fetchall()


def list_session_uploaded_documents(
    session_id: str, query: str = "", limit: int = 50
) -> dict[str, Any]:
    _validate_session_scope(session_id)
    if not isinstance(limit, int) or not 1 <= limit <= MAX_LINKED_UPLOADS_PER_SESSION:
        raise ValueError("limit must be an integer from 1 to 100")
    with _db() as conn:
        rows = _session_upload_rows(conn, session_id, query=query, limit=limit)
        documents = []
        for row in rows:
            status = _document_processing_snapshot_conn(conn, row["document_id"])
            documents.append(
                {
                    "session_id": session_id,
                    "document_id": row["document_id"],
                    "upload_id": row["upload_id"],
                    "upload_session_id": row["upload_session_id"],
                    "saved_path": row["relative_path"],
                    "name": row["name"],
                    "mime": row["mime_type"],
                    "size": row["byte_size"],
                    "sha256": row["sha256"],
                    "reused": bool(row["reused"]),
                    "linked_at": row["linked_at"],
                    "source_session_id": row["source_session_id"],
                    "created_at": row["upload_created_at"],
                    "processing": status,
                    "summary": (status["extraction"] or {}).get("summary")
                    if status["extraction"]
                    else "",
                    "classification": (status["extraction"] or {}).get("tags")
                    if status["extraction"]
                    else [],
                    "warnings": (status["extraction"] or {}).get("warnings")
                    if status["extraction"]
                    else [],
                    "preview_available": bool(
                        status["extraction"] and status["extraction"]["status"] == "complete"
                    ),
                }
            )
    return {"count": len(documents), "documents": documents}


def get_session_uploaded_document_status(session_id: str, document_id: str) -> dict[str, Any]:
    _validate_session_scope(session_id)
    if not DOCUMENT_ID_PATTERN.fullmatch(document_id or ""):
        raise ValueError("invalid document_id")
    with _db() as conn:
        row = conn.execute(
            "SELECT i.document_id, i.upload_id, i.relative_path, i.name, i.mime_type, i.byte_size, "
            "i.sha256, i.created_at AS upload_created_at, l.linked_at, l.source_session_id, "
            "(SELECT COUNT(*) FROM docops_upload_instances all_i WHERE all_i.document_id = i.document_id) > 1 AS reused "
            "FROM docops_chat_session_upload_instances l "
            "JOIN docops_upload_instances i ON i.instance_id = l.instance_id "
            "WHERE l.session_id=? AND i.document_id=? "
            "ORDER BY l.linked_at DESC LIMIT 1",
            (session_id, document_id),
        ).fetchone()
        if row is None:
            return {"error": f"document '{document_id}' is not linked to session '{session_id}'"}
        status = _document_processing_snapshot_conn(conn, document_id)
    return {
        "session_id": session_id,
        "document_id": row["document_id"],
        "upload_id": row["upload_id"],
        "saved_path": row["relative_path"],
        "name": row["name"],
        "mime": row["mime_type"],
        "size": row["byte_size"],
        "sha256": row["sha256"],
        "reused": bool(row["reused"]),
        "linked_at": row["linked_at"],
        "source_session_id": row["source_session_id"],
        "created_at": row["upload_created_at"],
        "processing": status,
        "summary": (status["extraction"] or {}).get("summary") if status["extraction"] else "",
        "classification": (status["extraction"] or {}).get("tags") if status["extraction"] else [],
        "warnings": (status["extraction"] or {}).get("warnings") if status["extraction"] else [],
        "preview_available": bool(
            status["extraction"] and status["extraction"]["status"] == "complete"
        ),
    }


def get_session_uploaded_document_preview(
    session_id: str, document_id: str, max_chars: int = 4000
) -> dict[str, Any]:
    _validate_session_scope(session_id)
    if not DOCUMENT_ID_PATTERN.fullmatch(document_id or ""):
        raise ValueError("invalid document_id")
    if not isinstance(max_chars, int) or not 100 <= max_chars <= 20000:
        raise ValueError("max_chars must be an integer from 100 to 20000")
    with _db() as conn:
        linked = conn.execute(
            "SELECT 1 FROM docops_chat_session_upload_documents WHERE session_id=? AND document_id=?",
            (session_id, document_id),
        ).fetchone()
        if linked is None:
            return {"error": f"document '{document_id}' is not linked to session '{session_id}'"}
        extraction = conn.execute(
            "SELECT normalized_text, markdown_text, status, updated_at FROM docops_upload_extractions "
            "WHERE document_id=? AND extractor_version=?",
            (document_id, CURRENT_EXTRACTOR_VERSION),
        ).fetchone()
    if extraction is None:
        return {
            "session_id": session_id,
            "document_id": document_id,
            "preview_available": False,
            "preview": "",
            "preview_truncated": False,
        }
    text = extraction["normalized_text"] or extraction["markdown_text"] or ""
    return {
        "session_id": session_id,
        "document_id": document_id,
        "status": extraction["status"],
        "updated_at": extraction["updated_at"],
        "preview_available": extraction["status"] == "complete",
        "preview": text[:max_chars],
        "preview_truncated": len(text) > max_chars,
    }


def get_session_upload_batch_status(session_id: str, limit: int = 50) -> dict[str, Any]:
    listing = list_session_uploaded_documents(session_id, limit=limit)
    return {
        "session_id": session_id,
        "count": listing["count"],
        "documents": [
            {
                "document_id": item["document_id"],
                "upload_id": item["upload_id"],
                "saved_path": item["saved_path"],
                "name": item["name"],
                "mime": item["mime"],
                "reused": item["reused"],
                "queue_state": item["processing"]["queue_state"],
                "summary": item["summary"],
                "classification": item["classification"],
                "warnings": item["warnings"],
                "preview_available": item["preview_available"],
                "linked_at": item["linked_at"],
                "created_at": item["created_at"],
            }
            for item in listing["documents"]
        ],
    }


def retry_uploaded_document_processing(
    *,
    session_id: str,
    document_id: str = "",
    upload_id: str = "",
    saved_path: str = "",
) -> dict[str, Any]:
    _validate_session_scope(session_id)
    now = _utc_now_sql()
    with _db() as conn:
        if document_id:
            if not DOCUMENT_ID_PATTERN.fullmatch(document_id):
                raise ValueError("invalid document_id")
            linked = conn.execute(
                "SELECT 1 FROM docops_chat_session_upload_documents "
                "WHERE session_id=? AND document_id=?",
                (session_id, document_id),
            ).fetchone()
        else:
            if not UPLOAD_ID_PATTERN.fullmatch(upload_id or ""):
                raise ValueError("invalid upload_id")
            linked = conn.execute(
                "SELECT i.document_id FROM docops_chat_session_upload_instances l "
                "JOIN docops_upload_instances i ON i.instance_id=l.instance_id "
                "WHERE l.session_id=? AND i.upload_id=? AND (? = '' OR i.relative_path=?) "
                "ORDER BY l.linked_at DESC LIMIT 1",
                (session_id, upload_id, saved_path, saved_path),
            ).fetchone()
            if linked is not None:
                document_id = linked["document_id"]
        if linked is None:
            return {"error": "uploaded document is not linked to this session"}

        extraction = conn.execute(
            "SELECT status FROM docops_upload_extractions WHERE document_id=? AND extractor_version=?",
            (document_id, CURRENT_EXTRACTOR_VERSION),
        ).fetchone()
        job = conn.execute(
            "SELECT job_id, status FROM docops_upload_jobs WHERE document_id=? AND extractor_version=?",
            (document_id, CURRENT_EXTRACTOR_VERSION),
        ).fetchone()
        if extraction and extraction["status"] == "complete":
            return {"document_id": document_id, "queue_state": "complete", "retried": False}
        if job is None:
            conn.execute(
                "INSERT INTO docops_upload_jobs "
                "(job_id, document_id, document_sha256, extractor_version, status, attempts, max_attempts, "
                "next_attempt_at, created_at, updated_at, lease_owner, lease_expires_at, heartbeat_at, last_error, completed_at) "
                "SELECT ?, d.document_id, d.sha256, ?, 'queued', 0, ?, ?, ?, ?, NULL, NULL, NULL, '', NULL "
                "FROM docops_upload_documents d WHERE d.document_id=?",
                (
                    _job_id(
                        conn.execute(
                            "SELECT sha256 FROM docops_upload_documents WHERE document_id=?",
                            (document_id,),
                        ).fetchone()["sha256"],
                        CURRENT_EXTRACTOR_VERSION,
                    ),
                    CURRENT_EXTRACTOR_VERSION,
                    DEFAULT_JOB_MAX_ATTEMPTS,
                    now,
                    now,
                    now,
                    document_id,
                ),
            )
            return {"document_id": document_id, "queue_state": "queued", "retried": True}
        if job["status"] in {"queued", "running", "complete"}:
            return {"document_id": document_id, "queue_state": job["status"], "retried": False}
        conn.execute(
            "UPDATE docops_upload_jobs SET status='queued', attempts=0, "
            "next_attempt_at=?, lease_owner=NULL, lease_expires_at=NULL, heartbeat_at=NULL, "
            "last_error='', completed_at=NULL, updated_at=? WHERE job_id=?",
            (now, now, job["job_id"]),
        )
    return {"document_id": document_id, "queue_state": "queued", "retried": True}


def link_upload_session_to_chat_session(
    *, chat_session_id: str, source_session_id: str, limit: int = 500
) -> dict[str, Any]:
    _validate_session_scope(chat_session_id)
    _validate_session_scope(source_session_id)
    if source_session_id != chat_session_id and not source_session_id.startswith("upload_"):
        return {
            "error": (
                "source_session_id must match the target chat session or be an anonymous "
                "upload_* session."
            )
        }
    if not isinstance(limit, int) or not 1 <= limit <= 1000:
        raise ValueError("limit must be an integer from 1 to 1000")
    linked_at = _utc_now_sql()
    with _db() as conn:
        if not _chat_session_exists_conn(conn, chat_session_id):
            return {"error": f"chat session '{chat_session_id}' was not found"}
        source_rows = conn.execute(
            "SELECT document_id, upload_id, relative_path FROM docops_upload_instances "
            "WHERE session_id=? ORDER BY created_at DESC LIMIT ?",
            (source_session_id, limit),
        ).fetchall()
        if not source_rows:
            return {
                "chat_session_id": chat_session_id,
                "source_session_id": source_session_id,
                "matched": 0,
                "linked": 0,
            }
        rows = [dict(row) for row in source_rows]
        linked = _link_upload_instances_to_chat_session_conn(
            conn,
            target_session_id=chat_session_id,
            source_session_id=source_session_id,
            rows=rows,
            linked_at=linked_at,
        )
    return {
        "chat_session_id": chat_session_id,
        "source_session_id": source_session_id,
        "matched": len(source_rows),
        "linked": linked,
    }


def attach_session_upload_metadata(
    sessions: list[dict[str, Any]],
    *,
    per_session_limit: int = MAX_LINKED_UPLOADS_IN_SESSION_DETAIL,
) -> list[dict[str, Any]]:
    if not sessions:
        return sessions
    with _db() as conn:
        for item in sessions:
            session_id = str(item.get("id") or item.get("session_id") or "")
            if not SESSION_ID_PATTERN.fullmatch(session_id):
                item["linked_uploads"] = []
                item["linked_upload_count"] = 0
                continue
            rows = _session_upload_rows(conn, session_id, query="", limit=per_session_limit)
            linked = []
            for row in rows:
                status = _document_processing_snapshot_conn(conn, row["document_id"])
                linked.append(
                    {
                        "document_id": row["document_id"],
                        "upload_id": row["upload_id"],
                        "saved_path": row["relative_path"],
                        "name": row["name"],
                        "mime": row["mime_type"],
                        "reused": bool(row["reused"]),
                        "queue_state": status["queue_state"],
                        "summary": (
                            status["extraction"]["summary"] if status["extraction"] else ""
                        ),
                        "classification": (
                            status["extraction"]["tags"] if status["extraction"] else []
                        ),
                        "warnings": (
                            status["extraction"]["warnings"] if status["extraction"] else []
                        ),
                        "preview_available": bool(
                            status["extraction"] and status["extraction"]["status"] == "complete"
                        ),
                        "linked_at": row["linked_at"],
                        "created_at": row["upload_created_at"],
                    }
                )
            count_row = conn.execute(
                "SELECT COUNT(*) FROM docops_chat_session_upload_instances WHERE session_id=?",
                (session_id,),
            ).fetchone()
            item["linked_uploads"] = linked
            item["linked_upload_count"] = int(count_row[0]) if count_row else 0
    return sessions


def search_uploaded_documents(session_id: str, query: str = "", limit: int = 50) -> dict[str, Any]:
    """Search uploaded documents linked to a chat session with bounded metadata."""
    return list_session_uploaded_documents(session_id, query=query, limit=limit)


def get_uploaded_document_processing_status(
    session_id: str,
    document_id: str = "",
    upload_id: str = "",
    saved_path: str = "",
) -> dict[str, Any]:
    """Get processing state for one linked uploaded document."""
    if document_id:
        return get_session_uploaded_document_status(session_id, document_id)
    _validate_session_scope(session_id)
    if not UPLOAD_ID_PATTERN.fullmatch(upload_id or ""):
        raise ValueError("invalid upload_id")
    with _db() as conn:
        row = conn.execute(
            "SELECT i.document_id FROM docops_chat_session_upload_instances l "
            "JOIN docops_upload_instances i ON i.instance_id=l.instance_id "
            "WHERE l.session_id=? AND i.upload_id=? AND (? = '' OR i.relative_path=?) "
            "ORDER BY l.linked_at DESC LIMIT 1",
            (session_id, upload_id, saved_path, saved_path),
        ).fetchone()
    if row is None:
        return {"error": "uploaded document is not linked to this session"}
    return get_session_uploaded_document_status(session_id, row["document_id"])


def upload_inventory_count() -> int:
    """Return the number of source documents in the upload registry."""
    with _db() as conn:
        return conn.execute("SELECT COUNT(*) FROM docops_upload_instances").fetchone()[0]


def _enqueue_processing_job_conn(
    conn: sqlite3.Connection,
    *,
    document_id: str,
    document_sha256: str,
    extractor_version: str,
) -> None:
    now = _utc_now_sql()
    conn.execute(
        "INSERT INTO docops_upload_jobs "
        "(job_id, document_id, document_sha256, extractor_version, status, attempts, max_attempts, "
        "next_attempt_at, created_at, updated_at) VALUES (?,?,?,?,?,?,?,?,?,?) "
        "ON CONFLICT(document_sha256, extractor_version) DO NOTHING",
        (
            _job_id(document_sha256, extractor_version),
            document_id,
            document_sha256,
            extractor_version,
            "queued",
            0,
            DEFAULT_JOB_MAX_ATTEMPTS,
            now,
            now,
            now,
        ),
    )


def claim_processing_job(worker_id: str, lease_seconds: int = DEFAULT_LEASE_SECONDS) -> dict | None:
    """Atomically claim one queued/retry/stale-running job."""
    if not worker_id:
        raise ValueError("worker_id is required")
    now = _utc_now_sql()
    lease_expires = _plus_seconds_sql(now, lease_seconds)
    with _db() as conn:
        conn.execute("BEGIN IMMEDIATE")
        row = conn.execute(
            "SELECT job_id, attempts, max_attempts FROM docops_upload_jobs "
            "WHERE ("
            " (status IN ('queued', 'retry') AND next_attempt_at <= ?) "
            " OR (status = 'running' AND lease_expires_at IS NOT NULL AND lease_expires_at <= ?)"
            ") "
            "AND attempts < max_attempts "
            "ORDER BY next_attempt_at ASC, created_at ASC LIMIT 1",
            (now, now),
        ).fetchone()
        if row is None:
            conn.commit()
            return None
        attempts = int(row["attempts"]) + 1
        conn.execute(
            "UPDATE docops_upload_jobs SET status='running', attempts=?, lease_owner=?, "
            "lease_expires_at=?, heartbeat_at=?, updated_at=?, last_error='' WHERE job_id=?",
            (attempts, worker_id, lease_expires, now, now, row["job_id"]),
        )
        claimed = conn.execute(
            "SELECT job_id, document_id, document_sha256, extractor_version, status, attempts, "
            "max_attempts, lease_owner, lease_expires_at, heartbeat_at, next_attempt_at, "
            "last_error, created_at, updated_at, completed_at FROM docops_upload_jobs WHERE job_id=?",
            (row["job_id"],),
        ).fetchone()
        conn.commit()
    return dict(claimed) if claimed else None


def heartbeat_processing_job(
    job_id: str, worker_id: str, lease_seconds: int = DEFAULT_LEASE_SECONDS
) -> bool:
    """Renew the lease for a running job owned by worker_id."""
    now = _utc_now_sql()
    lease_expires = _plus_seconds_sql(now, lease_seconds)
    with _db() as conn:
        updated = conn.execute(
            "UPDATE docops_upload_jobs SET heartbeat_at=?, lease_expires_at=?, updated_at=? "
            "WHERE job_id=? AND status='running' AND lease_owner=?",
            (now, lease_expires, now, job_id, worker_id),
        )
    return bool(updated.rowcount)


def _mark_job_complete(conn: sqlite3.Connection, *, job_id: str, worker_id: str) -> None:
    now = _utc_now_sql()
    conn.execute(
        "UPDATE docops_upload_jobs SET status='complete', completed_at=?, updated_at=?, "
        "lease_owner=?, lease_expires_at=NULL, heartbeat_at=? WHERE job_id=? AND lease_owner=?",
        (now, now, worker_id, now, job_id, worker_id),
    )


def _mark_job_failed(
    conn: sqlite3.Connection, *, job_id: str, worker_id: str, error: str, terminal: bool
) -> None:
    now = _utc_now_sql()
    status = "failed" if terminal else "retry"
    next_attempt = now if terminal else _plus_seconds_sql(now, DEFAULT_JOB_RETRY_SECONDS)
    conn.execute(
        "UPDATE docops_upload_jobs SET status=?, last_error=?, next_attempt_at=?, updated_at=?, "
        "lease_owner=NULL, lease_expires_at=NULL, heartbeat_at=? WHERE job_id=? AND lease_owner=?",
        (status, _bounded(error, limit=MAX_ERROR_CHARS), next_attempt, now, now, job_id, worker_id),
    )


def _write_text_atomic(path: Path, content: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    temporary: Path | None = None
    try:
        with tempfile.NamedTemporaryFile(
            mode="w",
            encoding="utf-8",
            dir=path.parent,
            prefix=f".{path.name}.",
            suffix=".tmp",
            delete=False,
        ) as handle:
            temporary = Path(handle.name)
            handle.write(content)
            handle.flush()
            os.fsync(handle.fileno())
        temporary.chmod(0o600)
        temporary.replace(path)
    finally:
        if temporary is not None:
            temporary.unlink(missing_ok=True)


def _summarize_text(text: str) -> str:
    chunks = [piece.strip() for piece in _SENTENCE_SPLIT_RE.split(text) if piece.strip()]
    if not chunks:
        return ""
    return " ".join(chunks[:4])[:2_000]


def _classify_text(
    extension: str, text: str, metadata: dict[str, Any], mime_type: str
) -> list[str]:
    lowered = text.lower()
    tags = {
        f"ext:{extension.lstrip('.')}",
        f"mime:{mime_type.split('/')[0] if '/' in mime_type else mime_type}",
    }
    if metadata.get("line_count", 0) > 200:
        tags.add("shape:long")
    elif metadata.get("line_count", 0) > 0:
        tags.add("shape:short")
    if any(token in lowered for token in ("invoice", "amount", "total", "balance")):
        tags.add("topic:finance")
    if any(token in lowered for token in ("meeting", "agenda", "action item", "decision")):
        tags.add("topic:meeting")
    if any(token in lowered for token in ("security", "threat", "vulnerability", "risk")):
        tags.add("topic:security")
    if any(token in lowered for token in ("customer", "client", "account")):
        tags.add("topic:customer")
    if metadata.get("table_rows", 0) > 0:
        tags.add("shape:tabular")
    return sorted(tags)


def _normalize_text(value: str) -> str:
    lines = value.replace("\r\n", "\n").replace("\r", "\n").split("\n")
    normalized = "\n".join(_REPEATED_WS_RE.sub(" ", line.rstrip()) for line in lines)
    if len(normalized) <= MAX_TEXT_CHARS:
        return normalized
    return normalized[:MAX_TEXT_CHARS]


def _extract_text_file(path: Path) -> tuple[str, dict[str, Any]]:
    content = path.read_text(encoding="utf-8", errors="replace")
    normalized = _normalize_text(content)
    line_count = normalized.count("\n") + (1 if normalized else 0)
    return normalized, {
        "line_count": line_count,
        "character_count": len(normalized),
    }


def _extract_json_file(path: Path) -> tuple[str, dict[str, Any]]:
    parsed = json.loads(path.read_text(encoding="utf-8"))
    pretty = json.dumps(parsed, ensure_ascii=False, sort_keys=True, indent=2)
    normalized = _normalize_text(pretty)
    node_count = 0
    stack = [parsed]
    while stack and node_count < 50_000:
        node = stack.pop()
        node_count += 1
        if isinstance(node, dict):
            stack.extend(node.values())
        elif isinstance(node, list):
            stack.extend(node)
    return normalized, {
        "node_count": node_count,
        "character_count": len(normalized),
    }


def _extract_xml_file(path: Path) -> tuple[str, dict[str, Any]]:
    parser = etree.XMLParser(resolve_entities=False, no_network=True, recover=True)
    tree = etree.parse(str(path), parser)
    root = tree.getroot()
    text = " ".join(piece.strip() for piece in root.itertext() if piece.strip())
    normalized = _normalize_text(text)
    return normalized, {
        "root_tag": root.tag,
        "element_count": len(list(root.iter())),
        "character_count": len(normalized),
    }


def _extract_delimited_file(path: Path, delimiter: str) -> tuple[str, dict[str, Any]]:
    lines: list[str] = []
    row_count = 0
    max_columns = 0
    with path.open(encoding="utf-8", errors="replace", newline="") as handle:
        reader = csv.reader(handle, delimiter=delimiter)
        for row in reader:
            row_count += 1
            max_columns = max(max_columns, len(row))
            if len(lines) < 400:
                lines.append(" | ".join(cell.strip() for cell in row))
            if row_count >= 5_000:
                break
    normalized = _normalize_text("\n".join(lines))
    return normalized, {
        "table_rows": row_count,
        "max_columns": max_columns,
        "character_count": len(normalized),
    }


def _extract_docx_file(path: Path) -> tuple[str, dict[str, Any]]:
    paragraph_count = 0
    runs: list[str] = []
    with zipfile.ZipFile(path) as archive:
        xml = archive.read("word/document.xml")
    parser = etree.XMLParser(resolve_entities=False, no_network=True)
    root = etree.fromstring(xml, parser=parser)
    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    for paragraph in root.findall(".//w:p", namespaces=ns):
        parts = [
            node.text
            for node in paragraph.findall(".//w:t", namespaces=ns)
            if node.text and node.text.strip()
        ]
        if parts:
            paragraph_count += 1
            runs.append("".join(parts))
        if len(runs) >= 3_000:
            break
    normalized = _normalize_text("\n".join(runs))
    return normalized, {
        "paragraph_count": paragraph_count,
        "character_count": len(normalized),
    }


def _extract_xlsx_file(path: Path) -> tuple[str, dict[str, Any]]:
    workbook = load_workbook(path, read_only=True, data_only=True)
    sheet_names = workbook.sheetnames[:6]
    lines: list[str] = []
    row_count = 0
    try:
        for sheet_name in sheet_names:
            sheet = workbook[sheet_name]
            lines.append(f"# Sheet: {sheet_name}")
            for row in sheet.iter_rows(min_row=1, max_row=300, max_col=24, values_only=True):
                cells = ["" if value is None else str(value) for value in row]
                if any(cells):
                    lines.append(" | ".join(cells))
                    row_count += 1
                if row_count >= 2_000:
                    break
            if row_count >= 2_000:
                break
    finally:
        workbook.close()
    normalized = _normalize_text("\n".join(lines))
    return normalized, {
        "sheet_count": len(workbook.sheetnames),
        "scanned_sheets": len(sheet_names),
        "table_rows": row_count,
        "character_count": len(normalized),
    }


def _extract_pptx_file(path: Path) -> tuple[str, dict[str, Any]]:
    deck = Presentation(str(path))
    lines: list[str] = []
    scanned_shapes = 0
    for index, slide in enumerate(deck.slides, start=1):
        lines.append(f"# Slide {index}")
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                lines.append(text.strip())
                scanned_shapes += 1
            if scanned_shapes >= 3_000:
                break
        if scanned_shapes >= 3_000:
            break
    normalized = _normalize_text("\n".join(lines))
    return normalized, {
        "slide_count": len(deck.slides),
        "text_shape_count": scanned_shapes,
        "character_count": len(normalized),
    }


def _extract_pdf_file(path: Path) -> tuple[str, dict[str, Any]]:
    if PdfReader is None:
        raise ValueError("pypdf is required for PDF extraction")
    try:
        reader = PdfReader(str(path))
    except Exception as exc:
        raise ValueError("invalid or unreadable pdf content") from exc
    pages: list[str] = []
    extracted_pages = 0
    for page in reader.pages[:200]:
        try:
            text = page.extract_text() or ""
        except Exception as exc:
            raise ValueError("pdf page text extraction failed") from exc
        if text.strip():
            pages.append(text.strip())
        extracted_pages += 1
        if sum(len(piece) for piece in pages) >= MAX_TEXT_CHARS:
            break
    normalized = _normalize_text("\n\n".join(pages))
    return normalized, {
        "page_count": len(reader.pages),
        "scanned_pages": extracted_pages,
        "character_count": len(normalized),
    }


def _extract_rtf_file(path: Path) -> tuple[str, dict[str, Any]]:
    content = path.read_text(encoding="utf-8", errors="replace")
    content = _RTF_HEX_RE.sub("", content)
    content = _RTF_CONTROL_RE.sub(" ", content)
    content = content.replace("{", " ").replace("}", " ")
    normalized = _normalize_text(content)
    return normalized, {
        "character_count": len(normalized),
    }


def _libreoffice_extract(path: Path, extension: str) -> ExtractionResult:
    binary = shutil.which("libreoffice") or shutil.which("soffice")
    if not binary:
        warning = (
            f"LibreOffice is unavailable; install libreoffice to process '{extension}' documents."
        )
        return ExtractionResult(
            status="unavailable",
            markdown="",
            text="",
            metadata={"adapter": "libreoffice", "extension": extension},
            warnings=[warning],
            provenance={"adapter": "libreoffice", "binary": None},
            tags=[f"ext:{extension.lstrip('.')}"],
            summary="",
            error=warning,
        )
    with tempfile.TemporaryDirectory(prefix="pj-upload-lo-") as temp_dir:
        temp_root = Path(temp_dir)
        profile_dir = temp_root / "profile"
        profile_dir.mkdir(parents=True, exist_ok=True)
        output_dir = temp_root / "out"
        output_dir.mkdir(parents=True, exist_ok=True)
        env = dict(os.environ)
        env["HOME"] = str(temp_root)
        env["TMPDIR"] = str(temp_root)
        env["SAL_USE_VCLPLUGIN"] = "headless"
        env["SAL_DISABLE_SYNCHRONOUS_PRINTER_DETECTION"] = "1"
        cmd = [
            binary,
            "--headless",
            "--nologo",
            "--nodefault",
            "--nofirststartwizard",
            "--norestore",
            "--nolockcheck",
            "--safe-mode",
            f"-env:UserInstallation=file://{profile_dir.resolve().as_posix()}",
            "--convert-to",
            "txt:Text",
            "--outdir",
            str(output_dir),
            str(path),
        ]
        try:
            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=20,
                env=env,
                check=False,
            )
        except subprocess.TimeoutExpired:
            message = f"LibreOffice timed out while extracting '{path.name}'."
            return ExtractionResult(
                status="error",
                markdown="",
                text="",
                metadata={"adapter": "libreoffice", "extension": extension},
                warnings=[message],
                provenance={"adapter": "libreoffice", "binary": binary},
                tags=[f"ext:{extension.lstrip('.')}"],
                summary="",
                error=message,
            )
        txt_path = output_dir / f"{path.stem}.txt"
        if result.returncode != 0 or not txt_path.is_file():
            stderr = _bounded(
                result.stderr or result.stdout or "conversion failed", limit=MAX_ERROR_CHARS
            )
            message = f"LibreOffice conversion failed: {stderr}"
            return ExtractionResult(
                status="error",
                markdown="",
                text="",
                metadata={"adapter": "libreoffice", "extension": extension},
                warnings=[message],
                provenance={"adapter": "libreoffice", "binary": binary},
                tags=[f"ext:{extension.lstrip('.')}"],
                summary="",
                error=message,
            )
        extracted = _normalize_text(txt_path.read_text(encoding="utf-8", errors="replace"))
        metadata = {
            "adapter": "libreoffice",
            "extension": extension,
            "character_count": len(extracted),
            "stdout": _bounded(result.stdout, limit=MAX_ERROR_CHARS),
        }
        tags = _classify_text(extension, extracted, metadata, "application/octet-stream")
        summary = _summarize_text(extracted)
        markdown = _normalized_markdown(path.name, extracted, metadata, tags, summary)
        return ExtractionResult(
            status="complete",
            markdown=markdown,
            text=extracted,
            metadata=metadata,
            warnings=[],
            provenance={"adapter": "libreoffice", "binary": binary},
            tags=tags,
            summary=summary,
        )


def _normalized_markdown(
    name: str, text: str, metadata: dict[str, Any], tags: list[str], summary: str
) -> str:
    tag_line = ", ".join(tags)
    summary_block = summary if summary else "No summary available."
    excerpt = text[:MAX_MARKDOWN_CHARS]
    return textwrap.dedent(
        f"""\
        # {name}

        ## Summary

        {summary_block}

        ## Tags

        {tag_line}

        ## Metadata

        ```json
        {json.dumps(metadata, sort_keys=True, ensure_ascii=False, indent=2)}
        ```

        ## Extracted Text

        ```
        {excerpt}
        ```
        """
    ).strip()


def _extract_document(path: Path, extension: str, mime_type: str) -> ExtractionResult:
    warnings: list[str] = []
    metadata: dict[str, Any]
    extracted_text: str
    if extension in TEXT_EXTENSIONS:
        extracted_text, metadata = _extract_text_file(path)
    elif extension == ".json":
        extracted_text, metadata = _extract_json_file(path)
    elif extension in DELIMITED_EXTENSIONS:
        extracted_text, metadata = _extract_delimited_file(path, DELIMITED_EXTENSIONS[extension])
    elif extension == ".xml":
        extracted_text, metadata = _extract_xml_file(path)
    elif extension == ".docx":
        extracted_text, metadata = _extract_docx_file(path)
    elif extension == ".xlsx":
        extracted_text, metadata = _extract_xlsx_file(path)
    elif extension == ".pptx":
        extracted_text, metadata = _extract_pptx_file(path)
    elif extension == ".pdf":
        extracted_text, metadata = _extract_pdf_file(path)
    elif extension == ".rtf":
        extracted_text, metadata = _extract_rtf_file(path)
    elif extension in LIBREOFFICE_EXTENSIONS:
        return _libreoffice_extract(path, extension)
    else:
        message = f"No local extractor is available for extension '{extension}'."
        return ExtractionResult(
            status="unavailable",
            markdown="",
            text="",
            metadata={"extension": extension},
            warnings=[message],
            provenance={"adapter": "none"},
            tags=[f"ext:{extension.lstrip('.')}"] if extension else [],
            summary="",
            error=message,
        )

    tags = _classify_text(extension, extracted_text, metadata, mime_type)
    summary = _summarize_text(extracted_text)
    markdown = _normalized_markdown(path.name, extracted_text, metadata, tags, summary)
    return ExtractionResult(
        status="complete",
        markdown=markdown,
        text=extracted_text,
        metadata=metadata,
        warnings=warnings[:MAX_WARNINGS],
        provenance={"adapter": "python", "extension": extension, "mime_type": mime_type},
        tags=tags,
        summary=summary,
    )


def _store_derived_output(
    *,
    document_id: str,
    extractor_version: str,
    result: ExtractionResult,
) -> str:
    if not DOCUMENT_ID_PATTERN.fullmatch(document_id or ""):
        raise ValueError("invalid document_id")
    target_root = DERIVED_DIR / document_id / extractor_version
    target_root.mkdir(parents=True, exist_ok=True)
    _write_text_atomic(target_root / "normalized.md", result.markdown)
    _write_text_atomic(target_root / "normalized.txt", result.text)
    write_json_atomic(
        target_root / "metadata.json",
        {
            "status": result.status,
            "tags": result.tags,
            "summary": result.summary,
            "warnings": result.warnings,
            "provenance": result.provenance,
            "structural_metadata": result.metadata,
            "error": result.error or "",
            "extractor_version": extractor_version,
        },
    )
    return str(target_root)


def _persist_extraction_result(
    conn: sqlite3.Connection,
    *,
    document_id: str,
    extractor_version: str,
    result: ExtractionResult,
    derived_root: str,
) -> None:
    now = _utc_now_sql()
    conn.execute(
        "INSERT INTO docops_upload_extractions "
        "(document_id, extractor_version, status, markdown_text, normalized_text, "
        "structural_metadata_json, tags_json, summary_text, warnings_json, provenance_json, "
        "derived_root, error_detail, created_at, updated_at) VALUES (?,?,?,?,?,?,?,?,?,?,?,?,?,?) "
        "ON CONFLICT(document_id, extractor_version) DO UPDATE SET "
        "status=excluded.status, markdown_text=excluded.markdown_text, "
        "normalized_text=excluded.normalized_text, structural_metadata_json=excluded.structural_metadata_json, "
        "tags_json=excluded.tags_json, summary_text=excluded.summary_text, warnings_json=excluded.warnings_json, "
        "provenance_json=excluded.provenance_json, derived_root=excluded.derived_root, "
        "error_detail=excluded.error_detail, updated_at=excluded.updated_at",
        (
            document_id,
            extractor_version,
            result.status,
            result.markdown,
            result.text,
            json.dumps(result.metadata, sort_keys=True),
            json.dumps(result.tags, sort_keys=True),
            result.summary,
            json.dumps(result.warnings),
            json.dumps(result.provenance, sort_keys=True),
            derived_root,
            result.error or "",
            now,
            now,
        ),
    )


def _load_document_for_job(conn: sqlite3.Connection, document_id: str) -> sqlite3.Row:
    row = conn.execute(
        "SELECT d.document_id, d.sha256, d.byte_size, d.canonical_path, d.canonical_name, d.canonical_mime_type "
        "FROM docops_upload_documents d WHERE d.document_id=?",
        (document_id,),
    ).fetchone()
    if row is None:
        raise ValueError(f"unknown upload document '{document_id}'")
    return row


def _process_claimed_job(job: dict[str, Any], worker_id: str) -> dict[str, Any]:
    with _db() as conn:
        already_complete = conn.execute(
            "SELECT status FROM docops_upload_extractions WHERE document_id=? AND extractor_version=?",
            (job["document_id"], job["extractor_version"]),
        ).fetchone()
        if already_complete and already_complete["status"] == "complete":
            _mark_job_complete(conn, job_id=job["job_id"], worker_id=worker_id)
            return {"job_id": job["job_id"], "status": "complete", "reason": "already_complete"}

        doc = _load_document_for_job(conn, job["document_id"])
    path = Path(doc["canonical_path"])
    resolved = path.resolve(strict=True)
    resolved.relative_to(UPLOADS_DIR.resolve())
    if path.is_symlink() or not resolved.is_file():
        raise ValueError("document path is not a regular upload file")
    if resolved.stat().st_size != int(doc["byte_size"]):
        raise ValueError("document size mismatch during processing")
    if sha256_file(resolved) != doc["sha256"]:
        raise ValueError("document hash mismatch during processing")

    extension = resolved.suffix.lower()
    result = _extract_document(resolved, extension, doc["canonical_mime_type"])
    derived_root = _store_derived_output(
        document_id=doc["document_id"],
        extractor_version=job["extractor_version"],
        result=result,
    )
    with _db() as conn:
        _persist_extraction_result(
            conn,
            document_id=doc["document_id"],
            extractor_version=job["extractor_version"],
            result=result,
            derived_root=derived_root,
        )
        if result.status == "complete":
            _mark_job_complete(conn, job_id=job["job_id"], worker_id=worker_id)
            return {"job_id": job["job_id"], "status": "complete", "result_status": result.status}
        terminal = result.status == "unavailable" or job["attempts"] >= job["max_attempts"]
        _mark_job_failed(
            conn,
            job_id=job["job_id"],
            worker_id=worker_id,
            error=result.error or "extraction failed",
            terminal=terminal,
        )
        return {
            "job_id": job["job_id"],
            "status": "failed" if terminal else "retry",
            "result_status": result.status,
        }


def process_next_upload_job(worker_id: str) -> dict[str, Any] | None:
    """Claim and process a single upload job, with retry-safe status transitions."""
    job = claim_processing_job(worker_id=worker_id)
    if job is None:
        return None
    heartbeat_processing_job(job["job_id"], worker_id)
    try:
        return _process_claimed_job(job, worker_id)
    except (
        OSError,
        ValueError,
        sqlite3.Error,
        zipfile.BadZipFile,
        KeyError,
        etree.XMLSyntaxError,
        json.JSONDecodeError,
    ) as exc:
        terminal = bool(job["attempts"] >= job["max_attempts"])
        with _db() as conn:
            _mark_job_failed(
                conn,
                job_id=job["job_id"],
                worker_id=worker_id,
                error=str(exc),
                terminal=terminal,
            )
        return {
            "job_id": job["job_id"],
            "status": "failed" if terminal else "retry",
            "error": _bounded(str(exc), limit=MAX_ERROR_CHARS),
        }


def run_upload_processor_once(worker_id: str, max_jobs: int = 100) -> dict[str, Any]:
    """Process ready jobs once and return aggregate counts."""
    processed = 0
    completed = 0
    retried = 0
    failed = 0
    for _ in range(max_jobs):
        result = process_next_upload_job(worker_id)
        if result is None:
            break
        processed += 1
        if result["status"] == "complete":
            completed += 1
        elif result["status"] == "retry":
            retried += 1
        else:
            failed += 1
    return {
        "processed": processed,
        "completed": completed,
        "retried": retried,
        "failed": failed,
    }


def run_upload_processor_watch(worker_id: str, poll_interval_seconds: float = 2.0) -> None:
    """Continuously process jobs in a durable polling loop."""
    if poll_interval_seconds <= 0:
        raise ValueError("poll_interval_seconds must be positive")
    while True:
        result = run_upload_processor_once(worker_id, max_jobs=100)
        if result["processed"] == 0:
            time.sleep(poll_interval_seconds)


def list_processing_jobs(limit: int = 100) -> dict[str, Any]:
    if not isinstance(limit, int) or not 1 <= limit <= 500:
        raise ValueError("limit must be an integer from 1 to 500")
    with _db() as conn:
        rows = conn.execute(
            "SELECT job_id, document_id, document_sha256, extractor_version, status, attempts, "
            "max_attempts, lease_owner, lease_expires_at, heartbeat_at, next_attempt_at, last_error, "
            "created_at, updated_at, completed_at FROM docops_upload_jobs "
            "ORDER BY created_at DESC LIMIT ?",
            (limit,),
        ).fetchall()
    return {"count": len(rows), "jobs": [dict(row) for row in rows]}


UPLOAD_SCHEMAS = [
    {
        "type": "function",
        "name": "list_uploaded_documents",
        "description": (
            "List user-uploaded source documents registered in DocOps, "
            "optionally filtered by session or keyword."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "session_id": {"type": "string"},
                "query": {"type": "string"},
                "limit": {"type": "integer", "minimum": 1, "maximum": 100},
            },
            "required": [],
        },
    },
    {
        "type": "function",
        "name": "search_uploaded_documents",
        "description": (
            "Search uploaded documents linked to a chat session, including "
            "local processing/classification status."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "session_id": {"type": "string"},
                "query": {"type": "string"},
                "limit": {"type": "integer", "minimum": 1, "maximum": 100},
            },
            "required": ["session_id"],
        },
    },
    {
        "type": "function",
        "name": "get_uploaded_document",
        "description": (
            "Read metadata for an uploaded source document. Text-based formats "
            "also return a bounded content preview."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "upload_id": {"type": "string"},
                "session_id": {"type": "string"},
                "saved_path": {
                    "type": "string",
                    "description": "Specific saved path when an upload contains multiple files",
                },
            },
            "required": ["upload_id"],
        },
    },
    {
        "type": "function",
        "name": "get_uploaded_document_processing_status",
        "description": (
            "Inspect local processing status, summary, classification, and warnings "
            "for a document linked to the chat session."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "session_id": {"type": "string"},
                "document_id": {"type": "string"},
                "upload_id": {"type": "string"},
                "saved_path": {"type": "string"},
            },
            "required": ["session_id"],
        },
    },
    {
        "type": "function",
        "name": "retry_uploaded_document_processing",
        "description": (
            "Retry eligible failed local-processing jobs for an uploaded "
            "document linked to the chat session."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "session_id": {"type": "string"},
                "document_id": {"type": "string"},
                "upload_id": {"type": "string"},
                "saved_path": {"type": "string"},
            },
            "required": ["session_id"],
        },
    },
]

UPLOAD_DISPATCH = {
    "list_uploaded_documents": list_uploaded_documents,
    "search_uploaded_documents": search_uploaded_documents,
    "get_uploaded_document": get_uploaded_document,
    "get_uploaded_document_processing_status": get_uploaded_document_processing_status,
    "retry_uploaded_document_processing": retry_uploaded_document_processing,
}
