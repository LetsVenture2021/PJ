"""Safe, non-executing inspectors for supported publication formats."""

from __future__ import annotations

import html
import re
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from .base import ArtifactAdapter, Inspection
from .model import CanonicalDocument


class MarkdownAdapter(ArtifactAdapter):
    format = "md"
    renderer_version = "pj-markdown-1"

    def inspect(self, path):
        model = CanonicalDocument.from_markdown(Path(path).read_text(encoding="utf-8"))
        return Inspection(
            text_items=list(model.reading_order),
            headings=[b.text for b in model.blocks if b.kind == "heading"],
            tables=[[list(row) for row in b.rows] for b in model.blocks if b.kind == "table"],
            links=[link for b in model.blocks for link in b.metadata.get("links", ())],
            checks={"utf8": True},
        )


class HTMLAdapter(ArtifactAdapter):
    format = "html"
    renderer_version = "lxml-6"

    def inspect(self, path):
        from lxml import html as lh

        raw = Path(path).read_text(encoding="utf-8")
        result = Inspection()
        doc = lh.fromstring(raw)
        language = (doc.get("lang") or "").strip()
        result.checks.update(
            {
                "language": bool(re.fullmatch(r"[A-Za-z]{2,3}(?:-[A-Za-z0-9]+)*", language)),
                "charset": bool(doc.xpath("//meta[translate(@charset,'UTF','utf')='utf-8']")),
                "title": bool("".join(doc.xpath("//title/text()")).strip()),
                "landmarks": bool(doc.xpath("//main"))
                and bool(doc.xpath("//header|//nav"))
                and bool(doc.xpath("//footer")),
                "print_styles": "@media print" in raw,
                "responsive_overflow": bool(
                    re.search(r"(?:overflow|overflow-wrap|word-wrap|max-width)", raw)
                ),
            }
        )
        for name, passed in result.checks.items():
            if not passed:
                result.errors.append(
                    {"code": f"html_{name}", "message": f"HTML requirement failed: {name}"}
                )
        headings = doc.xpath("//h1|//h2|//h3|//h4|//h5|//h6")
        levels = [int(node.tag[1]) for node in headings]
        if any(current > previous + 1 for previous, current in zip(levels, levels[1:])):
            result.errors.append(
                {"code": "html_heading_order", "message": "heading levels are not logical"}
            )
        result.headings = [" ".join(node.itertext()).strip() for node in headings]
        result.text_items = [
            " ".join(node.itertext()).strip()
            for node in doc.xpath("//h1|//h2|//h3|//h4|//h5|//h6|//p|//li|//th|//td|//pre")
        ]
        for table in doc.xpath("//table"):
            if not table.xpath(".//th"):
                result.errors.append(
                    {"code": "html_table_headers", "message": "table has no header cells"}
                )
            result.tables.append(
                [
                    [" ".join(cell.itertext()).strip() for cell in row.xpath("./th|./td")]
                    for row in table.xpath(".//tr")
                ]
            )
        for anchor in doc.xpath("//a[@href]"):
            label = " ".join(anchor.itertext()).strip()
            target = anchor.get("href", "")
            result.links.append((label, target))
            if not label or label.casefold() in {"click here", "here", "link", "read more"}:
                result.errors.append(
                    {"code": "html_link_text", "message": "link text is empty or ambiguous"}
                )
        if any(not (img.get("alt") or "").strip() for img in doc.xpath("//img")):
            result.errors.append({"code": "html_alt_text", "message": "image is missing alt text"})
        if doc.xpath("//script|//*[@onload or @onclick or @onerror or @onmouseover]") or doc.xpath(
            "//*[@src[starts-with(.,'http://') or starts-with(.,'https://')]]"
        ):
            result.errors.append(
                {
                    "code": "html_active_content",
                    "message": "remote active content or unsafe script/event attribute found",
                }
            )
        return result


class PDFAdapter(ArtifactAdapter):
    format = "pdf"
    renderer_version = "pypdf-6"

    def inspect(self, path):
        from pypdf import PdfReader

        reader = PdfReader(path)
        result = Inspection()
        pages = []
        if not reader.pages:
            result.errors.append({"code": "pdf_page_count", "message": "PDF contains no pages"})
        for number, page in enumerate(reader.pages, 1):
            text = page.extract_text() or ""
            pages.append(text)
            if not text.strip():
                result.errors.append(
                    {
                        "code": "pdf_blank_page",
                        "message": f"page {number} is blank or has no extractable text",
                    }
                )
            box = page.mediabox
            if float(box.width) <= 0 or float(box.height) <= 0:
                result.errors.append(
                    {"code": "pdf_page_bounds", "message": f"page {number} has invalid bounds"}
                )
        result.text_items = [
            line.strip() for page in pages for line in page.splitlines() if line.strip()
        ]
        metadata = reader.metadata or {}
        if not metadata.get("/Title"):
            result.errors.append(
                {"code": "pdf_metadata", "message": "PDF title metadata is missing"}
            )
        root = reader.trailer.get("/Root", {})
        result.checks = {
            "page_count": len(reader.pages),
            "extractable_text": bool(result.text_items),
            "metadata": dict(metadata),
            "page_number_continuity": all(f"Page {i}" in text for i, text in enumerate(pages, 1)),
            "tagged": "/StructTreeRoot" in root,
            "tagging_supported": False,
            "margins_and_clipping": "verified_by_layout_renderer",
        }
        if not result.checks["page_number_continuity"]:
            result.errors.append(
                {"code": "pdf_page_numbers", "message": "page-number continuity failed"}
            )
        return result


class OOXMLAdapter(ArtifactAdapter):
    relationship_ns = {"r": "http://schemas.openxmlformats.org/package/2006/relationships"}

    def package(self, path, result):
        try:
            archive = zipfile.ZipFile(path)
            if archive.testzip():
                result.errors.append({"code": "ooxml_corrupt", "message": "corrupt OOXML member"})
            names = set(archive.namelist())
            macro = [
                name
                for name in names
                if "vbaproject" in name.casefold()
                or name.casefold().startswith(("word/activex/", "ppt/activex/", "xl/activex/"))
            ]
            external = []
            for name in names:
                if name.endswith(".rels"):
                    root = ET.fromstring(archive.read(name))
                    external += [
                        rel.get("Target", "") for rel in root if rel.get("TargetMode") == "External"
                    ]
            if macro:
                result.errors.append(
                    {"code": "ooxml_macro", "message": "macro or binary parts are prohibited"}
                )
            if external:
                result.errors.append(
                    {
                        "code": "ooxml_external_relationship",
                        "message": "external OOXML relationships are prohibited",
                    }
                )
            result.checks.update(
                {"package_valid": True, "macro_parts": macro, "external_relationships": external}
            )
            return archive, names
        except (zipfile.BadZipFile, ET.ParseError) as exc:
            result.errors.append(
                {"code": "ooxml_package", "message": f"invalid OOXML package: {exc}"}
            )
            return None, set()


class DOCXAdapter(OOXMLAdapter):
    format = "docx"
    renderer_version = "ooxml-safe-1"

    def inspect(self, path):
        result = Inspection()
        archive, names = self.package(path, result)
        if archive is None:
            return result
        ns = {
            "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
            "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        }
        required = {"word/document.xml", "docProps/core.xml"}
        if not required <= names:
            result.errors.append(
                {
                    "code": "docx_parts",
                    "message": "required document or core-property part missing",
                }
            )
            return result
        root = ET.fromstring(archive.read("word/document.xml"))
        order = []
        for element in root.findall(".//w:body/*", ns):
            if element.tag.endswith("}p"):
                text = "".join(element.itertext()).strip()
                style = element.find("./w:pPr/w:pStyle", ns)
                if text:
                    order.append(text)
                if style is not None and (style.get(f"{{{ns['w']}}}val") or "").startswith(
                    "Heading"
                ):
                    result.headings.append(text)
            elif element.tag.endswith("}tbl"):
                table = [
                    ["".join(cell.itertext()).strip() for cell in row.findall("./w:tc", ns)]
                    for row in element.findall("./w:tr", ns)
                ]
                result.tables.append(table)
                order.extend(cell for row in table for cell in row)
        result.text_items = order
        result.links = [
            (node.text or "", node.get(f"{{{ns['r']}}}id", ""))
            for node in root.findall(".//w:hyperlink", ns)
        ]
        result.checks.update(
            {
                "styles": "word/styles.xml" in names,
                "core_properties": True,
                "embedded_media": sorted(n for n in names if n.startswith("word/media/")),
                "table_semantics": True,
            }
        )
        archive.close()
        return result


class RTFAdapter(ArtifactAdapter):
    format = "rtf"
    renderer_version = "rtf-structural-1"

    def inspect(self, path):
        raw = Path(path).read_bytes().decode("latin-1")
        result = Inspection()
        if not raw.startswith("{\\rtf1") or raw.count("{") != raw.count("}"):
            result.errors.append(
                {"code": "rtf_structure", "message": "invalid RTF header or unbalanced groups"}
            )
        if re.search(r"\\object|\\objdata|\\field\s*{[^}]*DDE", raw, re.I):
            result.errors.append(
                {"code": "rtf_unsafe_object", "message": "unsafe embedded RTF object found"}
            )

        def unicode_char(match):
            value = int(match.group(1))
            return chr(value if value >= 0 else value + 65536)

        text = re.sub(r"\\u(-?\d+)\??", unicode_char, raw)
        text = re.sub(r"\\(?:par|line)\b ?", "\n", text)
        text = re.sub(r"\\'[0-9a-fA-F]{2}", " ", text)
        text = re.sub(r"\\[a-zA-Z]+-?\d* ?", " ", text)
        text = text.translate(str.maketrans({"{": " ", "}": " ", "\\": " "}))
        result.text_items = [
            html.unescape(value).strip() for value in re.split(r"[\r\n]+", text) if value.strip()
        ]
        result.checks = {
            "header": raw.startswith("{\\rtf1"),
            "balanced_groups": raw.count("{") == raw.count("}"),
            "unicode_controls": len(re.findall(r"\\u-?\d+", raw)),
        }
        return result


class XLSXAdapter(OOXMLAdapter):
    format = "xlsx"
    renderer_version = "openpyxl-3"

    def inspect(self, path):
        from openpyxl import load_workbook

        result = Inspection()
        archive, _ = self.package(path, result)
        if archive:
            archive.close()
        workbook = load_workbook(path, data_only=False, read_only=False)
        result.checks.update(
            {
                "sheet_names": workbook.sheetnames,
                "accessibility_title": workbook.properties.title or "",
                "hidden_content": [],
            }
        )
        for sheet in workbook.worksheets:
            if sheet.sheet_state != "visible":
                result.checks["hidden_content"].append(sheet.title)
            values = []
            for row in sheet.iter_rows():
                values.append(["" if cell.value is None else str(cell.value) for cell in row])
                for cell in row:
                    if cell.data_type == "f":
                        result.errors.append(
                            {
                                "code": "xlsx_formula_injection",
                                "message": f"formula found in {sheet.title}!{cell.coordinate}",
                            }
                        )
            result.tables.append(values)
            result.text_items.extend(cell for row in values for cell in row if cell)
            result.checks[sheet.title] = {
                "frozen_headers": str(sheet.freeze_panes or ""),
                "filters": sheet.auto_filter.ref or "",
                "print_area": str(sheet.print_area or ""),
                "print_titles": sheet.print_title_rows or "",
                "cell_types": sorted({cell.data_type for row in sheet.iter_rows() for cell in row}),
                "date_formats": sorted(
                    {
                        cell.number_format
                        for row in sheet.iter_rows()
                        for cell in row
                        if getattr(cell, "is_date", False)
                    }
                ),
            }
        if result.checks["hidden_content"]:
            result.errors.append(
                {"code": "xlsx_hidden_content", "message": "hidden sheets are prohibited"}
            )
        workbook.close()
        return result


class PPTXAdapter(OOXMLAdapter):
    format = "pptx"
    renderer_version = "python-pptx-1"

    def inspect(self, path):
        from pptx import Presentation

        result = Inspection()
        archive, _ = self.package(path, result)
        if archive:
            archive.close()
        deck = Presentation(path)
        for slide_number, slide in enumerate(deck.slides, 1):
            previous_top = -1
            for shape in slide.shapes:
                if (
                    shape.left < 0
                    or shape.top < 0
                    or shape.left + shape.width > deck.slide_width
                    or shape.top + shape.height > deck.slide_height
                ):
                    result.errors.append(
                        {
                            "code": "pptx_bounds",
                            "message": f"slide {slide_number} shape outside bounds",
                        }
                    )
                if shape.top < previous_top:
                    result.warnings.append(
                        {
                            "code": "pptx_reading_order",
                            "message": f"slide {slide_number} XML order differs from visual top-to-bottom order",
                        }
                    )
                previous_top = shape.top
                if getattr(shape, "has_text_frame", False):
                    text = shape.text.strip()
                    if text:
                        result.text_items.append(text)
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.size and run.font.size.pt < 7:
                                result.errors.append(
                                    {
                                        "code": "pptx_font_size",
                                        "message": f"slide {slide_number} text below the 7pt governed minimum",
                                    }
                                )
                if getattr(shape, "has_table", False):
                    table = [[cell.text for cell in row.cells] for row in shape.table.rows]
                    result.tables.append(table)
                    result.text_items.append(" ".join(cell for row in table for cell in row))
                if getattr(shape, "has_chart", False):
                    series_values = [
                        (series.name, list(series.values)) for series in shape.chart.series
                    ]
                    result.checks.setdefault("charts", []).append(
                        {"slide": slide_number, "series": series_values}
                    )
                    result.text_items.extend(
                        f"{name} " + " ".join(str(value) for value in values)
                        for name, values in series_values
                    )
                # python-pptx exposes decorative/name metadata, but not a stable alt-text API.
                if shape.shape_type in {13, 3} and not (shape.name or "").strip():
                    result.errors.append(
                        {
                            "code": "pptx_alt_text",
                            "message": f"slide {slide_number} visual lacks accessible description",
                        }
                    )
            notes = slide.notes_slide.notes_text_frame.text.strip()
            result.checks.setdefault("speaker_notes", []).append(notes)
            if notes:
                result.text_items.append(notes)
        result.checks.update(
            {
                "slide_count": len(deck.slides),
                "deterministic_rendering": True,
                "contrast": "renderer_brand_palette",
                "minimum_font_size_pt": 7,
            }
        )
        return result


ADAPTERS = {
    adapter.format: adapter
    for adapter in (
        MarkdownAdapter(),
        HTMLAdapter(),
        PDFAdapter(),
        DOCXAdapter(),
        RTFAdapter(),
        XLSXAdapter(),
        PPTXAdapter(),
    )
}
