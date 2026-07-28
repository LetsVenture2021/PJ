"""Structured, branded PowerPoint rendering for governed DocOps presentations."""
from __future__ import annotations

import io
import json
import math
import re
import tempfile
import textwrap
import zipfile
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

SLIDE_WIDTH = 13.333
SLIDE_HEIGHT = 7.5
MAX_SLIDES = 30
MAX_BULLETS = 7
MAX_TEXT_CHARS = 1300

BRAND = {
    "primary": "1557FF",
    "accent": "00DEEB",
    "ink": "050713",
    "muted": "526078",
    "paper": "F7F9FC",
    "line": "D8DCE6",
    "white": "FFFFFF",
    "success": "168A64",
    "warning": "D97706",
    "danger": "C73A5B",
}

SUPPORTED_LAYOUTS = {
    "title",
    "hero",
    "bullets",
    "two_column",
    "comparison",
    "metrics",
    "process",
    "cards",
    "table",
    "risk_matrix",
    "timeline",
    "bar_chart",
    "sources",
    "closing",
}

SLIDE_KEYS = {
    "layout",
    "title",
    "subtitle",
    "statement",
    "bullets",
    "left_title",
    "left_items",
    "right_title",
    "right_items",
    "metrics",
    "steps",
    "cards",
    "table",
    "chart",
    "notes",
    "sources",
    "eyebrow",
}

BLOCKING_MARKERS = ("[TBD", "[VERIFY CURRENT]", "{{", "TODO:")


class PresentationValidationError(ValueError):
    """Raised when a presentation specification cannot be rendered safely."""


def _clean_text(value, *, field: str, limit: int = 4000,
                required: bool = False) -> str:
    if value is None:
        value = ""
    if not isinstance(value, str):
        raise PresentationValidationError(f"{field} must be a string")
    text = value.strip()
    if required and not text:
        raise PresentationValidationError(f"{field} is required")
    if len(text) > limit:
        raise PresentationValidationError(
            f"{field} exceeds the {limit}-character limit"
        )
    return text


def _clean_string_list(value, *, field: str, limit: int = MAX_BULLETS,
                       item_limit: int = 240) -> list[str]:
    if value in (None, ""):
        return []
    if not isinstance(value, list):
        raise PresentationValidationError(f"{field} must be an array")
    if len(value) > limit:
        raise PresentationValidationError(
            f"{field} exceeds the {limit}-item limit"
        )
    return [
        _clean_text(item, field=f"{field}[{index}]", limit=item_limit, required=True)
        for index, item in enumerate(value)
    ]


def _clean_named_items(value, *, field: str, limit: int,
                       value_key: str = "body") -> list[dict]:
    if value in (None, ""):
        return []
    if not isinstance(value, list) or len(value) > limit:
        raise PresentationValidationError(
            f"{field} must be an array with at most {limit} entries"
        )
    cleaned = []
    for index, item in enumerate(value):
        if not isinstance(item, dict):
            raise PresentationValidationError(f"{field}[{index}] must be an object")
        cleaned.append({
            "title": _clean_text(
                item.get("title"),
                field=f"{field}[{index}].title",
                limit=100,
                required=True,
            ),
            value_key: _clean_text(
                item.get(value_key),
                field=f"{field}[{index}].{value_key}",
                limit=300,
                required=True,
            ),
        })
    return cleaned


def _clean_table(value, *, field: str) -> dict:
    if value in (None, "", {}):
        return {}
    if not isinstance(value, dict):
        raise PresentationValidationError(f"{field} must be an object")
    columns = _clean_string_list(
        value.get("columns"), field=f"{field}.columns", limit=5, item_limit=80
    )
    rows = value.get("rows") or []
    if not columns or not isinstance(rows, list) or not rows:
        raise PresentationValidationError(
            f"{field} requires non-empty columns and rows"
        )
    if len(rows) > 7:
        raise PresentationValidationError(f"{field}.rows exceeds the 7-row limit")
    cleaned_rows = []
    for row_index, row in enumerate(rows):
        if not isinstance(row, list) or len(row) != len(columns):
            raise PresentationValidationError(
                f"{field}.rows[{row_index}] must match the column count"
            )
        cleaned_rows.append([
            _clean_text(
                cell,
                field=f"{field}.rows[{row_index}][{col_index}]",
                limit=180,
                required=True,
            )
            for col_index, cell in enumerate(row)
        ])
    return {"columns": columns, "rows": cleaned_rows}


def _clean_chart(value, *, field: str) -> dict:
    if value in (None, "", {}):
        return {}
    if not isinstance(value, dict):
        raise PresentationValidationError(f"{field} must be an object")
    categories = _clean_string_list(
        value.get("categories"),
        field=f"{field}.categories",
        limit=8,
        item_limit=60,
    )
    series = value.get("series") or []
    if not categories or not isinstance(series, list) or not series:
        raise PresentationValidationError(
            f"{field} requires categories and at least one series"
        )
    if len(series) > 3:
        raise PresentationValidationError(f"{field}.series exceeds 3 entries")
    cleaned_series = []
    for index, item in enumerate(series):
        if not isinstance(item, dict):
            raise PresentationValidationError(
                f"{field}.series[{index}] must be an object"
            )
        values = item.get("values")
        if (
            not isinstance(values, list)
            or len(values) != len(categories)
            or any(
                isinstance(number, bool)
                or not isinstance(number, (int, float))
                or not math.isfinite(float(number))
                for number in values
            )
        ):
            raise PresentationValidationError(
                f"{field}.series[{index}].values must be numeric and match categories"
            )
        cleaned_series.append({
            "name": _clean_text(
                item.get("name"),
                field=f"{field}.series[{index}].name",
                limit=80,
                required=True,
            ),
            "values": [float(number) for number in values],
        })
    return {"categories": categories, "series": cleaned_series}


def normalize_spec(spec_json: str | dict) -> dict:
    """Parse and validate the authoritative presentation specification."""
    if isinstance(spec_json, str):
        try:
            raw = json.loads(spec_json)
        except json.JSONDecodeError as exc:
            raise PresentationValidationError(
                f"slides_json must be valid JSON: {exc}"
            ) from exc
    elif isinstance(spec_json, dict):
        raw = spec_json
    else:
        raise PresentationValidationError("presentation specification must be an object")
    if not isinstance(raw, dict):
        raise PresentationValidationError("presentation specification must be an object")

    raw_slides = raw.get("slides")
    if not isinstance(raw_slides, list) or not raw_slides:
        raise PresentationValidationError("slides must be a non-empty array")
    if len(raw_slides) > MAX_SLIDES:
        raise PresentationValidationError(
            f"slides exceeds the {MAX_SLIDES}-slide limit"
        )

    slides = []
    for index, raw_slide in enumerate(raw_slides):
        if not isinstance(raw_slide, dict):
            raise PresentationValidationError(f"slides[{index}] must be an object")
        unknown = sorted(set(raw_slide) - SLIDE_KEYS)
        if unknown:
            raise PresentationValidationError(
                f"slides[{index}] contains unknown fields: {unknown}"
            )
        layout = _clean_text(
            raw_slide.get("layout", "bullets"),
            field=f"slides[{index}].layout",
            limit=30,
            required=True,
        )
        if layout not in SUPPORTED_LAYOUTS:
            raise PresentationValidationError(
                f"slides[{index}].layout must be one of {sorted(SUPPORTED_LAYOUTS)}"
            )
        slide = {
            "layout": layout,
            "title": _clean_text(
                raw_slide.get("title"),
                field=f"slides[{index}].title",
                limit=140,
                required=True,
            ),
            "subtitle": _clean_text(
                raw_slide.get("subtitle"),
                field=f"slides[{index}].subtitle",
                limit=360,
            ),
            "statement": _clean_text(
                raw_slide.get("statement"),
                field=f"slides[{index}].statement",
                limit=500,
            ),
            "eyebrow": _clean_text(
                raw_slide.get("eyebrow"),
                field=f"slides[{index}].eyebrow",
                limit=60,
            ),
            "bullets": _clean_string_list(
                raw_slide.get("bullets"), field=f"slides[{index}].bullets"
            ),
            "left_title": _clean_text(
                raw_slide.get("left_title"),
                field=f"slides[{index}].left_title",
                limit=100,
            ),
            "left_items": _clean_string_list(
                raw_slide.get("left_items"),
                field=f"slides[{index}].left_items",
                limit=6,
            ),
            "right_title": _clean_text(
                raw_slide.get("right_title"),
                field=f"slides[{index}].right_title",
                limit=100,
            ),
            "right_items": _clean_string_list(
                raw_slide.get("right_items"),
                field=f"slides[{index}].right_items",
                limit=6,
            ),
            "metrics": _clean_named_items(
                raw_slide.get("metrics"),
                field=f"slides[{index}].metrics",
                limit=6,
                value_key="value",
            ),
            "steps": _clean_named_items(
                raw_slide.get("steps"),
                field=f"slides[{index}].steps",
                limit=7,
            ),
            "cards": _clean_named_items(
                raw_slide.get("cards"),
                field=f"slides[{index}].cards",
                limit=6,
            ),
            "table": _clean_table(
                raw_slide.get("table"), field=f"slides[{index}].table"
            ),
            "chart": _clean_chart(
                raw_slide.get("chart"), field=f"slides[{index}].chart"
            ),
            "notes": _clean_text(
                raw_slide.get("notes"),
                field=f"slides[{index}].notes",
                limit=3000,
            ),
            "sources": _clean_string_list(
                raw_slide.get("sources"),
                field=f"slides[{index}].sources",
                limit=8,
                item_limit=300,
            ),
        }
        _validate_layout_content(slide, index)
        serialized = json.dumps(slide, ensure_ascii=True)
        if len(serialized) > MAX_TEXT_CHARS * 3:
            raise PresentationValidationError(
                f"slides[{index}] exceeds the content-density limit"
            )
        for marker in BLOCKING_MARKERS:
            if marker in serialized:
                raise PresentationValidationError(
                    f"slides[{index}] contains unresolved marker {marker!r}"
                )
        slides.append(slide)

    return {
        "schema_version": "1.0",
        "title": _clean_text(
            raw.get("title"), field="title", limit=160, required=True
        ),
        "subtitle": _clean_text(raw.get("subtitle"), field="subtitle", limit=360),
        "audience": _clean_text(
            raw.get("audience"), field="audience", limit=160, required=True
        ),
        "brand": "aimhi",
        "slides": slides,
    }


def _validate_layout_content(slide: dict, index: int):
    layout = slide["layout"]
    has_body = bool(
        slide["subtitle"]
        or slide["statement"]
        or slide["bullets"]
        or slide["left_items"]
        or slide["right_items"]
        or slide["metrics"]
        or slide["steps"]
        or slide["cards"]
        or slide["table"]
        or slide["chart"]
        or slide["sources"]
        or layout in {"title", "closing"}
    )
    if not has_body:
        raise PresentationValidationError(
            f"slides[{index}] has no content for layout {layout!r}"
        )
    requirements = {
        "metrics": ("metrics",),
        "process": ("steps",),
        "timeline": ("steps",),
        "cards": ("cards",),
        "table": ("table",),
        "risk_matrix": ("table",),
        "bar_chart": ("chart",),
        "comparison": ("left_items", "right_items"),
        "two_column": ("left_items", "right_items"),
        "sources": ("sources",),
    }
    missing = [field for field in requirements.get(layout, ()) if not slide[field]]
    if missing:
        raise PresentationValidationError(
            f"slides[{index}] layout {layout!r} requires {missing}"
        )
    if (
        layout in {"metrics", "bar_chart"}
        and not slide["sources"]
        and "illustrative" not in (
            f"{slide['title']} {slide['notes']}".casefold()
        )
    ):
        raise PresentationValidationError(
            f"slides[{index}] quantitative content requires sources or "
            "an illustrative-data note"
        )
    density_limits = {
        "hero": 650,
        "bullets": 950,
        "two_column": 1000,
        "comparison": 1000,
        "metrics": 720,
        "process": 950,
        "cards": 1100,
        "table": 1200,
        "risk_matrix": 1200,
        "timeline": 1300,
        "bar_chart": 700,
        "sources": 1300,
    }
    visible_text = [
        slide["title"],
        slide["subtitle"],
        slide["statement"],
        *slide["bullets"],
        *slide["left_items"],
        *slide["right_items"],
        *slide["sources"],
    ]
    visible_text.extend(
        f"{item['title']} {item.get('value', item.get('body', ''))}"
        for field in ("metrics", "steps", "cards")
        for item in slide[field]
    )
    visible_text.extend(
        str(cell)
        for row in slide["table"].get("rows", [])
        for cell in row
    )
    visible_text.extend(slide["table"].get("columns", []))
    visible_chars = sum(len(text) for text in visible_text)
    limit = density_limits.get(layout, 800)
    if visible_chars > limit:
        raise PresentationValidationError(
            f"slides[{index}] exceeds the {layout} layout density limit "
            f"({visible_chars}>{limit} visible characters)"
        )
    if layout == "process" and any(
        len(item["body"]) > 180 for item in slide["steps"]
    ):
        raise PresentationValidationError(
            f"slides[{index}] process step body exceeds 180 characters"
        )
    if layout == "cards" and any(
        len(item["body"]) > 240 for item in slide["cards"]
    ):
        raise PresentationValidationError(
            f"slides[{index}] card body exceeds 240 characters"
        )


def spec_to_markdown(spec: dict) -> str:
    """Render a human-reviewable companion without becoming the PPTX source."""
    normalized = normalize_spec(spec)
    lines = [
        f"# {normalized['title']}",
        "",
        f"**Audience:** {normalized['audience']}",
        "",
    ]
    if normalized["subtitle"]:
        lines.extend([normalized["subtitle"], ""])
    for index, slide in enumerate(normalized["slides"], start=1):
        lines.extend([
            f"## Slide {index}: {slide['title']}",
            "",
            f"**Layout:** {slide['layout']}",
            "",
        ])
        if slide["subtitle"]:
            lines.extend([slide["subtitle"], ""])
        if slide["statement"]:
            lines.extend([f"> {slide['statement']}", ""])
        for bullet in slide["bullets"]:
            lines.append(f"- {bullet}")
        for field, heading in (
            ("left_items", slide["left_title"] or "Left"),
            ("right_items", slide["right_title"] or "Right"),
        ):
            if slide[field]:
                lines.extend(["", f"### {heading}"])
                lines.extend(f"- {item}" for item in slide[field])
        for field in ("metrics", "steps", "cards"):
            if slide[field]:
                lines.extend(["", f"### {field.replace('_', ' ').title()}"])
                lines.extend(
                    f"- **{item['title']}:** "
                    f"{item.get('value', item.get('body', ''))}"
                    for item in slide[field]
                )
        if slide["table"]:
            columns = slide["table"]["columns"]
            lines.extend([
                "",
                "### Table",
                "",
                "| "
                + " | ".join(
                    str(cell).replace("|", "\\|") for cell in columns
                )
                + " |",
                "| " + " | ".join("---" for _ in columns) + " |",
            ])
            lines.extend(
                "| "
                + " | ".join(
                    str(cell).replace("|", "\\|") for cell in row
                )
                + " |"
                for row in slide["table"]["rows"]
            )
        if slide["chart"]:
            categories = slide["chart"]["categories"]
            series = slide["chart"]["series"]
            lines.extend([
                "",
                "### Chart data",
                "",
                "| Category | "
                + " | ".join(
                    item["name"].replace("|", "\\|") for item in series
                )
                + " |",
                "| --- | " + " | ".join("---" for _ in series) + " |",
            ])
            for category_index, category in enumerate(categories):
                lines.append(
                    "| "
                    + str(category).replace("|", "\\|")
                    + " | "
                    + " | ".join(
                        f"{item['values'][category_index]:g}"
                        for item in series
                    )
                    + " |"
                )
        if slide["sources"]:
            lines.extend(["", "### Sources"])
            lines.extend(f"- {source}" for source in slide["sources"])
        if slide["notes"]:
            lines.extend(["", "### Speaker notes", slide["notes"]])
        lines.append("")
    return "\n".join(lines).strip() + "\n"


def _rgb(hex_value: str) -> RGBColor:
    value = hex_value.lstrip("#")
    return RGBColor.from_string(value)


def _set_background(slide, color: str):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color)


def _add_textbox(slide, x, y, width, height, text, *, size=24,
                 color="ink", bold=False, align=PP_ALIGN.LEFT,
                 font_name="Aptos", valign=MSO_ANCHOR.TOP,
                 margin=0.05):
    box = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(width), Inches(height)
    )
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(BRAND.get(color, color))
    return box


def _add_rect(slide, x, y, width, height, *, fill="paper", line="line",
              radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(BRAND.get(fill, fill))
    shape.line.color.rgb = _rgb(BRAND.get(line, line))
    return shape


def _add_brand_mark(slide, *, dark=False):
    color = "white" if dark else "ink"
    _add_textbox(
        slide, 0.55, 0.28, 2.3, 0.32, "AIMHI / PJ",
        size=10, color=color, bold=True,
    )
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.68), Inches(0.72), Inches(0.04)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = _rgb(BRAND["accent"])
    accent.line.fill.background()


def _add_slide_title(slide, slide_data):
    _add_brand_mark(slide)
    if slide_data["eyebrow"]:
        _add_textbox(
            slide, 0.65, 0.85, 11.8, 0.3, slide_data["eyebrow"].upper(),
            size=10, color="primary", bold=True,
        )
    _add_textbox(
        slide, 0.65, 1.12, 12.0, 0.8, slide_data["title"],
        size=28, color="ink", bold=True,
    )
    if slide_data["subtitle"]:
        _add_textbox(
            slide, 0.68, 1.86, 11.7, 0.5, slide_data["subtitle"],
            size=13, color="muted",
        )


def _add_footer(slide, doc_id: str, version: int, index: int, total: int,
                *, dark=False):
    color = "white" if dark else "muted"
    _add_textbox(
        slide, 0.65, 7.06, 8.0, 0.22,
        f"{doc_id} v{version} | INTERNAL",
        size=8, color=color,
    )
    _add_textbox(
        slide, 11.8, 7.06, 0.85, 0.22, f"{index}/{total}",
        size=8, color=color, align=PP_ALIGN.RIGHT,
    )


def _add_draft_watermark(slide):
    watermark = _add_textbox(
        slide, 3.35, 3.05, 6.6, 0.9, "DRAFT",
        size=48, color="line", bold=True, align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    watermark.rotation = -24


def _add_bullets(slide, items, x, y, width, height, *, size=19,
                 color="ink"):
    box = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(width), Inches(height)
    )
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.05)
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = _rgb(BRAND.get(color, color))
        paragraph.space_after = Pt(10)
        paragraph.line_spacing = 1.08
    return box


def _render_title(slide, data):
    _set_background(slide, BRAND["ink"])
    _add_brand_mark(slide, dark=True)
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.16), Inches(7.5)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = _rgb(BRAND["accent"])
    accent.line.fill.background()
    _add_textbox(
        slide, 0.82, 1.65, 11.6, 1.5, data["title"],
        size=38, color="white", bold=True,
    )
    subtitle = data["subtitle"] or data["statement"]
    if subtitle:
        _add_textbox(
            slide, 0.85, 3.35, 10.9, 1.1, subtitle,
            size=20, color="white",
        )


def _render_hero(slide, data):
    _set_background(slide, BRAND["paper"])
    _add_slide_title(slide, data)
    _add_rect(slide, 0.65, 2.6, 12.0, 3.45, fill="white", line="line")
    statement = data["statement"] or data["subtitle"]
    _add_textbox(
        slide, 1.05, 3.12, 11.15, 1.75, statement,
        size=27, color="ink", bold=True, valign=MSO_ANCHOR.MIDDLE,
    )
    if data["bullets"]:
        _add_bullets(slide, data["bullets"], 1.08, 5.1, 10.9, 0.72, size=14)


def _render_bullets(slide, data):
    _set_background(slide, BRAND["white"])
    _add_slide_title(slide, data)
    _add_bullets(slide, data["bullets"], 0.8, 2.55, 7.55, 3.95, size=19)
    _add_rect(slide, 8.75, 2.55, 3.85, 3.45, fill="primary", line="primary")
    statement = data["statement"] or data["subtitle"] or data["title"]
    _add_textbox(
        slide, 9.1, 2.95, 3.15, 2.7, statement,
        size=21, color="white", bold=True, valign=MSO_ANCHOR.MIDDLE,
    )


def _render_columns(slide, data, *, comparison=False):
    _set_background(slide, BRAND["white"])
    _add_slide_title(slide, data)
    left_fill = "paper"
    right_fill = "primary" if comparison else "paper"
    _add_rect(slide, 0.65, 2.45, 5.92, 4.05, fill=left_fill, line="line")
    _add_rect(
        slide, 6.78, 2.45, 5.88, 4.05,
        fill=right_fill, line="primary" if comparison else "line",
    )
    _add_textbox(
        slide, 0.98, 2.75, 5.2, 0.45,
        data["left_title"] or "Current state",
        size=18, color="ink", bold=True,
    )
    _add_textbox(
        slide, 7.12, 2.75, 5.1, 0.45,
        data["right_title"] or "PJ",
        size=18, color="white" if comparison else "ink", bold=True,
    )
    _add_bullets(
        slide, data["left_items"], 0.98, 3.35, 5.15, 2.85, size=15
    )
    _add_bullets(
        slide, data["right_items"], 7.12, 3.35, 5.08, 2.85, size=15,
        color="white" if comparison else "ink",
    )


def _render_metrics(slide, data):
    _set_background(slide, BRAND["white"])
    _add_slide_title(slide, data)
    metrics = data["metrics"]
    columns = 3
    rows = math.ceil(len(metrics) / columns)
    card_w = 3.75
    card_h = 1.55 if rows > 1 else 2.25
    for index, metric in enumerate(metrics):
        row, col = divmod(index, columns)
        x = 0.72 + col * 4.16
        y = 2.55 + row * (card_h + 0.35)
        _add_rect(slide, x, y, card_w, card_h, fill="paper", line="line")
        _add_textbox(
            slide, x + 0.25, y + 0.22, card_w - 0.5, 0.55,
            metric["value"], size=25, color="primary", bold=True,
        )
        _add_textbox(
            slide, x + 0.25, y + 0.83, card_w - 0.5, card_h - 0.95,
            metric["title"], size=13, color="muted", bold=True,
        )


def _render_process(slide, data, *, timeline=False):
    _set_background(slide, BRAND["white"])
    _add_slide_title(slide, data)
    steps = data["steps"]
    if timeline:
        gap = 0.12
        available_height = 4.25
        height = min(
            0.82,
            (available_height - gap * (len(steps) - 1)) / len(steps),
        )
        if height < 0.44:
            raise PresentationValidationError(
                "Timeline contains too many steps for the slide bounds"
            )
        for index, step in enumerate(steps):
            y = 2.45 + index * (height + gap)
            _add_rect(slide, 1.0, y, 0.65, height, fill="primary", line="primary")
            _add_textbox(
                slide, 1.0, y + 0.05, 0.65, height - 0.1, str(index + 1),
                size=15, color="white", bold=True, align=PP_ALIGN.CENTER,
                valign=MSO_ANCHOR.MIDDLE,
            )
            _add_textbox(
                slide, 1.95, y + 0.02, 2.4, height - 0.04,
                step["title"], size=15, color="ink", bold=True,
                valign=MSO_ANCHOR.MIDDLE,
            )
            _add_textbox(
                slide, 4.45, y + 0.02, 7.75, height - 0.04,
                step["body"], size=13, color="muted",
                valign=MSO_ANCHOR.MIDDLE,
            )
        return

    width = min(2.15, 11.7 / len(steps) - 0.12)
    total_width = len(steps) * width + (len(steps) - 1) * 0.15
    start_x = (SLIDE_WIDTH - total_width) / 2
    for index, step in enumerate(steps):
        x = start_x + index * (width + 0.15)
        _add_rect(slide, x, 2.7, width, 3.05, fill="paper", line="line")
        _add_textbox(
            slide, x + 0.15, 2.93, 0.45, 0.45, str(index + 1),
            size=18, color="primary", bold=True,
        )
        _add_textbox(
            slide, x + 0.15, 3.5, width - 0.3, 0.65,
            step["title"], size=15, color="ink", bold=True,
        )
        _add_textbox(
            slide, x + 0.15, 4.25, width - 0.3, 1.2,
            step["body"], size=11, color="muted",
        )


def _render_cards(slide, data):
    _set_background(slide, BRAND["white"])
    _add_slide_title(slide, data)
    cards = data["cards"]
    columns = 3
    rows = math.ceil(len(cards) / columns)
    card_h = 1.7 if rows > 1 else 3.2
    for index, card in enumerate(cards):
        row, col = divmod(index, columns)
        x = 0.72 + col * 4.15
        y = 2.45 + row * (card_h + 0.25)
        _add_rect(slide, x, y, 3.75, card_h, fill="paper", line="line")
        _add_textbox(
            slide, x + 0.25, y + 0.22, 3.25, 0.5,
            card["title"], size=16, color="primary", bold=True,
        )
        _add_textbox(
            slide, x + 0.25, y + 0.82, 3.25, card_h - 1.0,
            card["body"], size=12, color="muted",
        )


def _render_table(slide, data, *, risk=False):
    _set_background(slide, BRAND["white"])
    _add_slide_title(slide, data)
    table_data = data["table"]
    columns = table_data["columns"]
    rows = table_data["rows"]
    shape = slide.shapes.add_table(
        len(rows) + 1, len(columns),
        Inches(0.65), Inches(2.45), Inches(12.0), Inches(4.05),
    )
    table = shape.table
    for col_index, heading in enumerate(columns):
        cell = table.cell(0, col_index)
        cell.text = heading
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(
            BRAND["danger"] if risk else BRAND["primary"]
        )
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.color.rgb = _rgb(BRAND["white"])
            paragraph.font.bold = True
            paragraph.font.size = Pt(12)
    for row_index, row in enumerate(rows, start=1):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(
                BRAND["paper"] if row_index % 2 == 0 else BRAND["white"]
            )
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.color.rgb = _rgb(BRAND["ink"])
                paragraph.font.size = Pt(10.5)
    for col_index in range(len(columns)):
        table.columns[col_index].width = Inches(12.0 / len(columns))


def _render_chart(slide, data):
    _set_background(slide, BRAND["white"])
    _add_slide_title(slide, data)
    chart_data = ChartData()
    chart_data.categories = data["chart"]["categories"]
    for series in data["chart"]["series"]:
        chart_data.add_series(series["name"], series["values"])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.9), Inches(2.45), Inches(11.55), Inches(4.05),
        chart_data,
    ).chart
    chart.has_legend = len(data["chart"]["series"]) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.value_axis.has_major_gridlines = True
    chart.category_axis.tick_labels.font.size = Pt(11)
    chart.value_axis.tick_labels.font.size = Pt(10)


def _render_sources(slide, data):
    _set_background(slide, BRAND["paper"])
    _add_slide_title(slide, data)
    _add_bullets(slide, data["sources"], 0.85, 2.45, 11.7, 4.0, size=13)


def _render_closing(slide, data):
    _set_background(slide, BRAND["ink"])
    _add_brand_mark(slide, dark=True)
    _add_textbox(
        slide, 0.92, 1.05, 11.2, 0.45, data["title"].upper(),
        size=12, color="accent", bold=True,
    )
    statement = data["statement"] or data["subtitle"] or data["title"]
    _add_textbox(
        slide, 0.9, 1.85, 11.5, 2.0, statement,
        size=34, color="white", bold=True, valign=MSO_ANCHOR.MIDDLE,
    )
    if data["bullets"]:
        _add_bullets(slide, data["bullets"], 0.95, 4.3, 10.8, 1.5,
                     size=16, color="white")


def _expected_notes_text(slide: dict) -> str:
    notes = slide["notes"]
    if slide["sources"]:
        source_block = "Sources:\n" + "\n".join(
            f"- {source}" for source in slide["sources"]
        )
        notes = f"{notes}\n\n{source_block}".strip()
    return notes


def render_pptx(spec: dict, output_path: Path, *, doc_id: str,
                version: int, status: str, source_sha256: str) -> dict:
    """Render and validate a native PPTX package from a normalized spec."""
    normalized = normalize_spec(spec)
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_WIDTH)
    presentation.slide_height = Inches(SLIDE_HEIGHT)
    blank = presentation.slide_layouts[6]
    total = len(normalized["slides"])

    for index, data in enumerate(normalized["slides"], start=1):
        slide = presentation.slides.add_slide(blank)
        layout = data["layout"]
        if layout == "title":
            _render_title(slide, data)
        elif layout == "hero":
            _render_hero(slide, data)
        elif layout == "bullets":
            _render_bullets(slide, data)
        elif layout in {"two_column", "comparison"}:
            _render_columns(slide, data, comparison=layout == "comparison")
        elif layout == "metrics":
            _render_metrics(slide, data)
        elif layout in {"process", "timeline"}:
            _render_process(slide, data, timeline=layout == "timeline")
        elif layout == "cards":
            _render_cards(slide, data)
        elif layout in {"table", "risk_matrix"}:
            _render_table(slide, data, risk=layout == "risk_matrix")
        elif layout == "bar_chart":
            _render_chart(slide, data)
        elif layout == "sources":
            _render_sources(slide, data)
        elif layout == "closing":
            _render_closing(slide, data)
        else:
            raise PresentationValidationError(f"unsupported layout {layout!r}")

        dark = layout in {"title", "closing"}
        _add_footer(slide, doc_id, version, index, total, dark=dark)
        if status != "final":
            _add_draft_watermark(slide)
        slide.notes_slide.notes_text_frame.text = _expected_notes_text(data)

    properties = presentation.core_properties
    properties.title = normalized["title"]
    properties.subject = f"Audience: {normalized['audience']}"
    properties.author = "PJ / Aimhi"
    properties.keywords = (
        f"{doc_id}, version {version}, {status}, integrity {source_sha256}"
    )
    properties.comments = "Generated by governed DocOps presentation rendering."
    with tempfile.NamedTemporaryFile(
        dir=output_path.parent,
        prefix=f".{output_path.name}.",
        suffix=".tmp",
        delete=False,
    ) as handle:
        temporary = Path(handle.name)
    try:
        presentation.save(temporary)
        _canonicalize_pptx(temporary)
        validation = validate_pptx(temporary, normalized)
        temporary.replace(output_path)
    finally:
        temporary.unlink(missing_ok=True)
    return {
        "slides": total,
        "layouts": [slide["layout"] for slide in normalized["slides"]],
        "validation": validation,
    }


def _canonicalize_pptx(path: Path):
    """Rewrite ZIP metadata so identical specs produce identical package hashes."""
    path = Path(path)
    with tempfile.NamedTemporaryFile(
        dir=path.parent,
        prefix=f".{path.name}.",
        suffix=".canonical",
        delete=False,
    ) as handle:
        temporary = Path(handle.name)
    try:
        with zipfile.ZipFile(path, "r") as source:
            entries = [
                (
                    info,
                    _canonicalize_ooxml_member(
                        source.read(info.filename), info.filename
                    ),
                )
                for info in source.infolist()
            ]
        with zipfile.ZipFile(
            temporary, "w", compression=zipfile.ZIP_DEFLATED, compresslevel=9
        ) as target:
            for info, data in sorted(entries, key=lambda item: item[0].filename):
                canonical = zipfile.ZipInfo(
                    info.filename, date_time=(1980, 1, 1, 0, 0, 0)
                )
                canonical.compress_type = zipfile.ZIP_DEFLATED
                canonical.external_attr = info.external_attr
                canonical.internal_attr = info.internal_attr
                canonical.create_system = 0
                canonical.comment = info.comment
                target.writestr(canonical, data)
        temporary.replace(path)
    finally:
        temporary.unlink(missing_ok=True)


def _canonicalize_ooxml_member(data: bytes, filename: str) -> bytes:
    normalized_name = filename.replace("\\", "/").casefold()
    if normalized_name.endswith("docprops/core.xml"):
        for field in (b"created", b"modified"):
            pattern = (
                rb"(<dcterms:" + field + rb"\b[^>]*>)[^<]*"
                rb"(</dcterms:" + field + rb">)"
            )
            data = re.sub(
                pattern,
                rb"\g<1>1980-01-01T00:00:00Z\g<2>",
                data,
            )
    if normalized_name.endswith((".xlsx", ".docx", ".pptx")):
        return _canonicalize_embedded_package(data, filename)
    return data


def _canonicalize_embedded_package(data: bytes, filename: str) -> bytes:
    source_buffer = io.BytesIO(data)
    target_buffer = io.BytesIO()
    try:
        with zipfile.ZipFile(source_buffer, "r") as source:
            entries = [
                (
                    info,
                    _canonicalize_ooxml_member(
                        source.read(info.filename), info.filename
                    ),
                )
                for info in source.infolist()
            ]
        with zipfile.ZipFile(
            target_buffer, "w", compression=zipfile.ZIP_DEFLATED, compresslevel=9
        ) as target:
            for info, content in sorted(
                entries, key=lambda item: item[0].filename
            ):
                canonical = zipfile.ZipInfo(
                    info.filename, date_time=(1980, 1, 1, 0, 0, 0)
                )
                canonical.compress_type = zipfile.ZIP_DEFLATED
                canonical.external_attr = info.external_attr
                canonical.internal_attr = info.internal_attr
                canonical.create_system = 0
                canonical.comment = info.comment
                target.writestr(canonical, content)
    except zipfile.BadZipFile as exc:
        raise PresentationValidationError(
            f"embedded OOXML package {filename!r} is invalid"
        ) from exc
    return target_buffer.getvalue()


def _shape_texts(slide) -> list[str]:
    values = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = shape.text.strip()
            if text:
                values.append(text)
    return values


def validate_pptx(path: Path, expected_spec: dict) -> dict:
    """Reopen a PPTX and verify its governed package contract."""
    normalized = normalize_spec(expected_spec)
    path = Path(path)
    if not path.is_file() or path.stat().st_size < 10_000:
        raise PresentationValidationError("PPTX package is missing or unexpectedly small")
    try:
        with zipfile.ZipFile(path, "r") as package:
            names = set(package.namelist())
            required_parts = {
                "[Content_Types].xml",
                "_rels/.rels",
                "ppt/presentation.xml",
                "ppt/_rels/presentation.xml.rels",
            }
            if not required_parts.issubset(names):
                raise PresentationValidationError(
                    "PPTX package is missing required OOXML parts"
                )
            content_types = package.read("[Content_Types].xml")
            if (
                b"application/vnd.openxmlformats-officedocument."
                b"presentationml.presentation.main+xml"
                not in content_types
            ):
                raise PresentationValidationError(
                    "PPTX package has an invalid presentation content type"
                )
            if package.testzip() is not None:
                raise PresentationValidationError("PPTX package contains a corrupt ZIP member")
    except zipfile.BadZipFile as exc:
        raise PresentationValidationError("PPTX package is not a valid ZIP archive") from exc
    presentation = Presentation(path)
    if len(presentation.slides) != len(normalized["slides"]):
        raise PresentationValidationError("PPTX slide count does not match specification")
    if (
        round(presentation.slide_width / Inches(1), 3) != round(SLIDE_WIDTH, 3)
        or round(presentation.slide_height / Inches(1), 3) != round(SLIDE_HEIGHT, 3)
    ):
        raise PresentationValidationError("PPTX is not widescreen 16:9")
    for index, (slide, expected) in enumerate(
        zip(presentation.slides, normalized["slides"]), start=1
    ):
        for shape in slide.shapes:
            if (
                shape.left < 0
                or shape.top < 0
                or shape.left + shape.width > presentation.slide_width
                or shape.top + shape.height > presentation.slide_height
            ):
                raise PresentationValidationError(
                    f"slide {index} contains content outside slide bounds"
                )
        texts = _shape_texts(slide)
        expected_title = expected["title"].casefold()
        if not any(expected_title in text.casefold() for text in texts):
            raise PresentationValidationError(
                f"slide {index} does not contain its expected title"
            )
        combined = "\n".join(texts)
        if any(marker in combined for marker in BLOCKING_MARKERS):
            raise PresentationValidationError(
                f"slide {index} contains unresolved markers"
            )
        expected_notes = _expected_notes_text(expected)
        actual_notes = (
            slide.notes_slide.notes_text_frame.text
            .replace("\r\n", "\n")
            .replace("\r", "\n")
            .strip()
        )
        if actual_notes != expected_notes.strip():
            raise PresentationValidationError(
                f"slide {index} notes/sources do not match specification"
            )
        if expected["table"]:
            tables = [
                shape.table
                for shape in slide.shapes
                if getattr(shape, "has_table", False)
            ]
            if len(tables) != 1:
                raise PresentationValidationError(
                    f"slide {index} is missing its reviewable table"
                )
            actual_matrix = [
                [cell.text.strip() for cell in row.cells]
                for row in tables[0].rows
            ]
            expected_matrix = [
                [str(cell).strip() for cell in row]
                for row in (
                    [expected["table"]["columns"]]
                    + expected["table"]["rows"]
                )
            ]
            if actual_matrix != expected_matrix:
                raise PresentationValidationError(
                    f"slide {index} table content does not match specification"
                )
        if expected["chart"]:
            charts = [
                shape.chart
                for shape in slide.shapes
                if getattr(shape, "has_chart", False)
            ]
            if len(charts) != 1 or len(charts[0].series) != len(
                expected["chart"]["series"]
            ):
                raise PresentationValidationError(
                    f"slide {index} chart content does not match specification"
                )
            chart = charts[0]
            categories = [
                str(category.label)
                for category in chart.plots[0].categories
            ]
            if categories != expected["chart"]["categories"]:
                raise PresentationValidationError(
                    f"slide {index} chart categories do not match specification"
                )
            actual_series = [
                (series.name, tuple(float(value) for value in series.values))
                for series in charts[0].series
            ]
            expected_series = [
                (series["name"], tuple(float(value) for value in series["values"]))
                for series in expected["chart"]["series"]
            ]
            if actual_series != expected_series:
                raise PresentationValidationError(
                    f"slide {index} chart series do not match specification"
                )
    return {
        "status": "validated",
        "slide_count": len(presentation.slides),
        "widescreen": True,
        "package_bytes": path.stat().st_size,
        "mime_type": (
            "application/vnd.openxmlformats-officedocument."
            "presentationml.presentation"
        ),
    }


def render_previews(spec: dict, destination: Path) -> list[Path]:
    """Render deterministic preview images for visual review and overflow checks."""
    normalized = normalize_spec(spec)
    destination = Path(destination)
    destination.mkdir(parents=True, exist_ok=True)
    paths = []
    width, height = 1600, 900
    for index, slide in enumerate(normalized["slides"], start=1):
        dark = slide["layout"] in {"title", "closing"}
        background = BRAND["ink"] if dark else (
            BRAND["paper"] if slide["layout"] in {"hero", "sources"} else BRAND["white"]
        )
        image = Image.new("RGB", (width, height), f"#{background}")
        draw = ImageDraw.Draw(image)
        title_color = f"#{BRAND['white'] if dark else BRAND['ink']}"
        muted_color = f"#{BRAND['white'] if dark else BRAND['muted']}"
        draw.rectangle((0, 0, 20, height), fill=f"#{BRAND['accent']}")
        draw.text((70, 42), "AIMHI / PJ", fill=title_color)
        title = "\n".join(textwrap.wrap(slide["title"], width=38))
        draw.multiline_text(
            (75, 150 if dark else 115), title, fill=title_color,
            font=ImageFont.load_default(size=38), spacing=12,
        )
        body_values = []
        if slide["subtitle"]:
            body_values.append(slide["subtitle"])
        if slide["statement"]:
            body_values.append(slide["statement"])
        body_values.extend(f"• {value}" for value in slide["bullets"])
        for values, heading in (
            (slide["left_items"], slide["left_title"] or "Left"),
            (slide["right_items"], slide["right_title"] or "Right"),
        ):
            if values:
                body_values.append(
                    f"{heading}: " + "; ".join(str(value) for value in values)
                )
        for field in ("metrics", "steps", "cards"):
            if slide[field]:
                body_values.append(field.replace("_", " ").title() + ":")
                body_values.extend(
                    f"{item['title']}: {item.get('value', item.get('body', ''))}"
                    for item in slide[field]
                )
        if slide["table"]:
            body_values.append(" | ".join(slide["table"]["columns"]))
            body_values.extend(
                " | ".join(str(cell) for cell in row)
                for row in slide["table"]["rows"]
            )
        if slide["chart"]:
            for series in slide["chart"]["series"]:
                body_values.append(
                    f"{series['name']}: "
                    + ", ".join(
                        f"{category} {series['values'][category_index]:g}"
                        for category_index, category in enumerate(
                            slide["chart"]["categories"]
                        )
                    )
                )
        body_values.extend(
            f"Source: {source}" for source in slide["sources"]
        )
        if (
            slide["layout"] not in {"title", "closing"}
            and not any(str(value).strip() for value in body_values)
        ):
            raise PresentationValidationError(
                f"slide {index} preview has no substantive review content"
            )
        body = "\n\n".join(
            "\n".join(textwrap.wrap(value, width=72)) for value in body_values
        )
        body_y = 360 if dark else 285
        body_font = None
        body_spacing = 8
        for font_size in (22, 20, 18, 16, 14):
            candidate = ImageFont.load_default(size=font_size)
            bounds = draw.multiline_textbbox(
                (0, 0), body, font=candidate, spacing=body_spacing
            )
            if (
                bounds[2] - bounds[0] <= width - 160
                and bounds[3] - bounds[1] <= 820 - body_y
            ):
                body_font = candidate
                break
        if body and body_font is None:
            raise PresentationValidationError(
                f"slide {index} preview content exceeds the review canvas"
            )
        if body:
            draw.multiline_text(
                (80, body_y), body, fill=muted_color,
                font=body_font, spacing=body_spacing,
            )
        draw.text(
            (75, 845), f"{index}/{len(normalized['slides'])}",
            fill=muted_color,
        )
        path = destination / f"slide-{index:02d}.png"
        with tempfile.NamedTemporaryFile(
            dir=destination,
            prefix=f".{path.name}.",
            suffix=".tmp",
            delete=False,
        ) as handle:
            temporary = Path(handle.name)
        try:
            image.save(temporary, "PNG")
            if temporary.stat().st_size < 2_000:
                raise PresentationValidationError(
                    f"slide {index} preview was unexpectedly empty"
                )
            temporary.replace(path)
        finally:
            temporary.unlink(missing_ok=True)
        paths.append(path)
    return paths
