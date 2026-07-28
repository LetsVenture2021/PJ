import hashlib
import json
import tempfile
import threading
import time
import unittest
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path
from unittest.mock import patch

from pptx import Presentation

import docops
import presentationops


def presentation_fixture():
    return {
        "title": "PJ Internal Value Proposition",
        "subtitle": "Governed AI execution for internal stakeholders",
        "audience": "Internal stakeholders",
        "slides": [
            {
                "layout": "title",
                "title": "PJ: Governed AI Operations",
                "subtitle": "From intent to verified execution",
            },
            {
                "layout": "hero",
                "title": "Core proposition",
                "statement": "One operating layer for research, action, evidence, and deliverables.",
            },
            {
                "layout": "bullets",
                "title": "Operating problem",
                "bullets": ["Context is fragmented.", "Execution lacks receipts."],
                "statement": "PJ closes the intent-to-outcome gap.",
            },
            {
                "layout": "two_column",
                "title": "Connected workflow",
                "left_title": "Inputs",
                "left_items": ["Conversations", "Documents"],
                "right_title": "Outputs",
                "right_items": ["Decisions", "Artifacts"],
            },
            {
                "layout": "comparison",
                "title": "Why PJ is different",
                "left_title": "Conventional assistant",
                "left_items": ["Answer-centric", "Limited continuity"],
                "right_title": "PJ",
                "right_items": ["Outcome-centric", "Governed continuity"],
            },
            {
                "layout": "metrics",
                "title": "Verified production foundation",
                "metrics": [
                    {"title": "DocOps templates", "value": "437"},
                    {"title": "Template aliases", "value": "885"},
                    {"title": "Coding capabilities", "value": "8"},
                ],
                "notes": "Counts are sourced from the production inventory.",
                "sources": ["PJ capability snapshot, verified 2026-07-28"],
            },
            {
                "layout": "process",
                "title": "Execution path",
                "steps": [
                    {"title": "Aim", "body": "Define the outcome."},
                    {"title": "Act", "body": "Use governed tools."},
                    {"title": "Verify", "body": "Collect evidence."},
                ],
            },
            {
                "layout": "cards",
                "title": "High-value use cases",
                "cards": [
                    {"title": "Executive operations", "body": "Decisions and commitments."},
                    {"title": "Document operations", "body": "Versioned native artifacts."},
                    {"title": "Code operations", "body": "Read-only governed evidence."},
                ],
            },
            {
                "layout": "table",
                "title": "Capability matrix",
                "table": {
                    "columns": ["Capability", "State", "Control"],
                    "rows": [
                        ["Web research", "Active", "Citations"],
                        ["MCP", "Active", "Owner approval"],
                    ],
                },
            },
            {
                "layout": "risk_matrix",
                "title": "Risk controls",
                "table": {
                    "columns": ["Risk", "Control", "Evidence"],
                    "rows": [
                        ["Format substitution", "Artifact contract", "Package validation"],
                        ["Unverified claim", "Source notes", "Capability snapshot"],
                    ],
                },
            },
            {
                "layout": "timeline",
                "title": "Rollout",
                "steps": [
                    {"title": "Week 1", "body": "Select workflows."},
                    {"title": "Week 2", "body": "Run proof of value."},
                    {"title": "Week 3", "body": "Review evidence."},
                ],
            },
            {
                "layout": "bar_chart",
                "title": "Illustrative scorecard",
                "chart": {
                    "categories": ["Baseline", "Target"],
                    "series": [{"name": "Workflow coverage", "values": [1, 3]}],
                },
                "notes": "Targets are recommendations, not measured outcomes.",
            },
            {
                "layout": "sources",
                "title": "Sources and evidence",
                "sources": [
                    "PJ production capability snapshot",
                    "PJ public /health contract",
                ],
            },
            {
                "layout": "closing",
                "title": "Decision",
                "statement": "Authorize a focused, evidence-based proof of value.",
            },
        ],
    }


class TestPresentationOps(unittest.TestCase):
    def setUp(self):
        self.temp_dir = tempfile.TemporaryDirectory()
        self.root = Path(self.temp_dir.name)
        self.old_db = docops._DB_PATH
        self.old_docs = docops.DOCS_DIR
        self.old_exports = docops.EXPORTS_DIR
        docops._DB_PATH = self.root / "test.sqlite3"
        docops.DOCS_DIR = self.root / "documents"
        docops.DOCS_DIR.mkdir()
        docops.EXPORTS_DIR = docops.DOCS_DIR / "exports"
        docops.EXPORTS_DIR.mkdir()

    def tearDown(self):
        docops._DB_PATH = self.old_db
        docops.DOCS_DIR = self.old_docs
        docops.EXPORTS_DIR = self.old_exports
        self.temp_dir.cleanup()

    def test_all_core_layouts_render_with_notes_and_previews(self):
        spec = presentationops.normalize_spec(presentation_fixture())
        output = self.root / "all-layouts.pptx"
        result = presentationops.render_pptx(
            spec,
            output,
            doc_id="DOC-layouts",
            version=1,
            status="final",
            source_sha256="source-hash",
        )
        previews = presentationops.render_previews(
            spec, self.root / "previews"
        )

        self.assertEqual(result["slides"], len(spec["slides"]))
        self.assertEqual(len(previews), len(spec["slides"]))
        reopened = Presentation(output)
        self.assertEqual(len(reopened.slides), len(spec["slides"]))
        self.assertIn(
            "PJ capability snapshot",
            reopened.slides[5].notes_slide.notes_text_frame.text,
        )
        self.assertTrue(any(shape.has_chart for shape in reopened.slides[11].shapes))

    def test_governed_pptx_export_is_registered_and_idempotent(self):
        spec = presentation_fixture()
        drafted = docops.draft_presentation(
            spec["title"],
            spec["audience"],
            json.dumps(spec),
            subtitle=spec["subtitle"],
            finalize=True,
        )
        first = docops.export_document(drafted["doc_id"], "pptx")
        second = docops.export_document(drafted["doc_id"], "pptx")

        self.assertEqual(first["artifact"]["status"], "ready")
        self.assertTrue(first["audience_ready"])
        self.assertEqual(first["preview_count"], len(spec["slides"]))
        self.assertEqual(
            first["artifact"]["artifact_id"],
            second["artifact"]["artifact_id"],
        )
        self.assertEqual(
            first["artifact"]["sha256"],
            second["artifact"]["sha256"],
        )
        resolved = docops.resolve_export_artifact(
            first["artifact"]["artifact_id"]
        )
        self.assertEqual(resolved["sha256"], second["artifact"]["sha256"])
        self.assertNotIn("path", resolved)

    def test_artifact_id_is_immutable_when_draft_becomes_final(self):
        spec = presentation_fixture()
        drafted = docops.draft_presentation(
            spec["title"], spec["audience"], json.dumps(spec)
        )
        draft_export = docops.export_document(drafted["doc_id"], "pptx")
        finalized = docops.finalize_document(drafted["doc_id"])
        final_export = docops.export_document(drafted["doc_id"], "pptx")

        self.assertEqual(finalized["status"], "final")
        self.assertNotEqual(
            draft_export["artifact"]["artifact_id"],
            final_export["artifact"]["artifact_id"],
        )
        self.assertNotEqual(
            draft_export["artifact"]["sha256"],
            final_export["artifact"]["sha256"],
        )
        self.assertEqual(
            docops.resolve_export_artifact(
                draft_export["artifact"]["artifact_id"]
            )["sha256"],
            draft_export["artifact"]["sha256"],
        )
        self.assertEqual(
            docops.resolve_export_artifact(
                final_export["artifact"]["artifact_id"]
            )["sha256"],
            final_export["artifact"]["sha256"],
        )

    def test_companion_divergence_blocks_draft_and_final_presentations(self):
        for finalize in (False, True):
            with self.subTest(finalize=finalize):
                spec = presentation_fixture()
                drafted = docops.draft_presentation(
                    spec["title"],
                    spec["audience"],
                    json.dumps(spec),
                    finalize=finalize,
                )
                path = Path(drafted["path"])
                tampered = path.read_text().replace(
                    "One operating layer for research, action, evidence, and deliverables.",
                    "Companion text changed independently.",
                )
                path.write_text(tampered)
                with docops._db() as conn:
                    conn.execute(
                        "UPDATE docops_documents SET sha256=? "
                        "WHERE doc_id=? AND version=?",
                        (
                            hashlib.sha256(tampered.encode()).hexdigest(),
                            drafted["doc_id"],
                            drafted["version"],
                        ),
                    )

                result = docops.finalize_document(drafted["doc_id"])

                self.assertEqual(result["status"], "blocked")
                self.assertIn("companion diverges", result["reason"])

    def test_seven_step_timeline_stays_within_slide_bounds(self):
        spec = presentation_fixture()
        timeline = next(
            slide for slide in spec["slides"]
            if slide["layout"] == "timeline"
        )
        timeline["steps"] = [
            {
                "title": f"Step {index}",
                "body": f"Complete governed phase {index}.",
            }
            for index in range(1, 8)
        ]
        path = self.root / "seven-step-timeline.pptx"

        result = presentationops.render_pptx(
            spec,
            path,
            doc_id="DOC-TIMELINE",
            version=1,
            status="final",
            source_sha256="a" * 64,
        )

        self.assertEqual(result["slides"], len(spec["slides"]))
        presentationops.validate_pptx(path, spec)

    def test_markdown_companion_and_previews_include_table_chart_sources(self):
        spec = presentation_fixture()
        markdown = presentationops.spec_to_markdown(spec)
        self.assertIn("### Table", markdown)
        self.assertIn("| Capability | State | Control |", markdown)
        self.assertIn("### Chart data", markdown)
        self.assertIn("| Category | Workflow coverage |", markdown)
        self.assertIn("PJ capability snapshot", markdown)

        previews = presentationops.render_previews(
            spec, self.root / "substantive-previews"
        )
        self.assertEqual(len(previews), len(spec["slides"]))
        self.assertTrue(all(path.stat().st_size > 2_000 for path in previews))

    def test_validation_compares_notes_tables_and_chart_order_exactly(self):
        spec = presentation_fixture()
        output = self.root / "semantic-validation.pptx"
        presentationops.render_pptx(
            spec,
            output,
            doc_id="DOC-semantics",
            version=1,
            status="final",
            source_sha256="a" * 64,
        )

        wrong_notes = presentation_fixture()
        wrong_notes["slides"][5]["sources"] = ["Different source"]
        with self.assertRaisesRegex(
            presentationops.PresentationValidationError, "notes/sources"
        ):
            presentationops.validate_pptx(output, wrong_notes)

        wrong_table = presentation_fixture()
        wrong_table["slides"][8]["table"]["rows"].reverse()
        with self.assertRaisesRegex(
            presentationops.PresentationValidationError, "table content"
        ):
            presentationops.validate_pptx(output, wrong_table)

        wrong_categories = presentation_fixture()
        wrong_categories["slides"][11]["chart"]["categories"].reverse()
        wrong_categories["slides"][11]["chart"]["series"][0]["values"].reverse()
        with self.assertRaisesRegex(
            presentationops.PresentationValidationError, "chart categories"
        ):
            presentationops.validate_pptx(output, wrong_categories)

    def test_preview_rejects_content_that_cannot_fit_without_truncation(self):
        spec = presentation_fixture()
        spec["slides"][2]["bullets"] = [
            f"{index}: " + ("A" * 140)
            for index in range(6)
        ]
        with self.assertRaisesRegex(
            presentationops.PresentationValidationError, "review canvas"
        ):
            presentationops.render_previews(spec, self.root / "overflow-previews")

    def test_chart_export_is_deterministic_across_wall_clock_seconds(self):
        spec = presentation_fixture()
        first = self.root / "first.pptx"
        second = self.root / "second.pptx"
        kwargs = {
            "doc_id": "DOC-deterministic",
            "version": 1,
            "status": "final",
            "source_sha256": "a" * 64,
        }
        presentationops.render_pptx(spec, first, **kwargs)
        time.sleep(1.1)
        presentationops.render_pptx(spec, second, **kwargs)

        self.assertEqual(
            hashlib.sha256(first.read_bytes()).hexdigest(),
            hashlib.sha256(second.read_bytes()).hexdigest(),
        )

    def test_concurrent_exports_share_one_valid_registered_artifact(self):
        spec = presentation_fixture()
        drafted = docops.draft_presentation(
            spec["title"],
            spec["audience"],
            json.dumps(spec),
            finalize=True,
        )
        with ThreadPoolExecutor(max_workers=2) as executor:
            exports = list(executor.map(
                lambda _: docops.export_document(drafted["doc_id"], "pptx"),
                range(2),
            ))

        self.assertEqual(
            len({export["artifact"]["artifact_id"] for export in exports}), 1
        )
        self.assertEqual(
            len({export["artifact"]["sha256"] for export in exports}), 1
        )
        resolved = docops.resolve_export_artifact(
            exports[0]["artifact"]["artifact_id"]
        )
        self.assertEqual(resolved["status"], "ready")
        self.assertFalse(any(
            path.name.endswith((".tmp", ".canonical"))
            for path in docops.EXPORTS_DIR.iterdir()
        ))

    def test_concurrent_non_pptx_exports_use_isolated_html_inputs(self):
        drafted = docops.draft_document(
            "status_report",
            "Concurrent Export",
            json.dumps({
                "Period": "Q3",
                "Highlights": "Immutable delivery",
                "Metrics": "All formats",
                "Blockers": "None",
                "Next Period Plan": "Ship",
            }),
            finalize=True,
        )
        barrier = threading.Barrier(2)
        source_paths = []
        source_lock = threading.Lock()

        def convert(command, **_kwargs):
            source = Path(command[3])
            output = Path(command[5])
            with source_lock:
                source_paths.append(source)
            barrier.wait(timeout=5)
            output.write_bytes(source.read_bytes() + command[2].encode())
            return docops.subprocess.CompletedProcess(command, 0, "", "")

        with patch.object(docops.subprocess, "run", side_effect=convert):
            with ThreadPoolExecutor(max_workers=2) as executor:
                exports = list(executor.map(
                    lambda format_name: docops.export_document(
                        drafted["doc_id"], format_name
                    ),
                    ("docx", "rtf"),
                ))

        self.assertEqual(len(set(source_paths)), 2)
        self.assertTrue(all(path.name.startswith(".") for path in source_paths))
        self.assertFalse(
            (docops.EXPORTS_DIR / (
                f"{drafted['doc_id']}-concurrent-export-v1.html"
            )).exists()
        )
        self.assertTrue(all(
            export["artifact"]["status"] == "ready" for export in exports
        ))

    def test_draft_is_watermarked_and_corruption_blocks_download(self):
        spec = presentation_fixture()
        drafted = docops.draft_presentation(
            spec["title"], spec["audience"], json.dumps(spec)
        )
        exported = docops.export_document(drafted["doc_id"], "pptx")
        reopened = Presentation(exported["path"])
        slide_text = "\n".join(
            shape.text
            for shape in reopened.slides[0].shapes
            if getattr(shape, "has_text_frame", False)
        )
        self.assertIn("DRAFT", slide_text)
        self.assertFalse(exported["audience_ready"])

        registered = docops.resolve_export_artifact(
            exported["artifact"]["artifact_id"], include_path=True
        )
        with Path(registered["path"]).open("ab") as handle:
            handle.write(b"corruption")
        resolved = docops.resolve_export_artifact(
            exported["artifact"]["artifact_id"]
        )
        self.assertEqual(resolved["status"], "blocked")
        self.assertEqual(resolved["error"], "artifact integrity mismatch")

    def test_legacy_document_can_receive_structured_v2_without_editing_v1(self):
        legacy = docops.draft_document(
            "status_report",
            "Legacy Presentation",
            json.dumps({
                "Period": "Q3",
                "Highlights": "Original v1",
                "Metrics": "N/A",
                "Blockers": "None",
                "Next Period Plan": "Create native PowerPoint",
            }),
        )
        v1_path = Path(legacy["path"])
        v1_before = v1_path.read_bytes()
        revised = docops.revise_presentation(
            legacy["doc_id"],
            json.dumps(presentation_fixture()),
            audience="Internal stakeholders",
            change_note="Corrected native presentation",
            finalize=True,
        )

        self.assertEqual(revised["version"], 2)
        self.assertEqual(v1_path.read_bytes(), v1_before)
        exported = docops.export_document(legacy["doc_id"], "pptx", version=2)
        self.assertEqual(exported["format"], "pptx")
        v1_export = docops.export_document(legacy["doc_id"], "pptx", version=1)
        self.assertEqual(v1_export["status"], "rejected")

    def test_strict_spec_rejects_unknown_fields_and_excess_density(self):
        invalid = presentation_fixture()
        invalid["slides"][0]["unsupported"] = True
        with self.assertRaises(presentationops.PresentationValidationError):
            presentationops.normalize_spec(invalid)

        invalid = presentation_fixture()
        invalid["slides"][11]["chart"]["series"][0]["values"][0] = float("nan")
        with self.assertRaises(presentationops.PresentationValidationError):
            presentationops.normalize_spec(invalid)

        invalid = presentation_fixture()
        invalid["slides"][2]["bullets"] = [str(index) for index in range(8)]
        with self.assertRaises(presentationops.PresentationValidationError):
            presentationops.normalize_spec(invalid)


if __name__ == "__main__":
    unittest.main()
