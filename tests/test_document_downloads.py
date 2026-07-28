import json
import os
import shutil
import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch

import chatlog
import docops
import realtime_server
from openpyxl import load_workbook
from pptx import Presentation


class TestDocumentDownloads(unittest.TestCase):
    def setUp(self):
        self.temp_dir = tempfile.TemporaryDirectory()
        root = Path(self.temp_dir.name)
        self.old_docops = {
            "_DB_PATH": docops._DB_PATH,
            "DOCS_DIR": docops.DOCS_DIR,
            "EXPORTS_DIR": docops.EXPORTS_DIR,
            "ARTIFACTS_DIR": docops.ARTIFACTS_DIR,
        }
        self.old_chat_db = chatlog._DB_PATH
        docops._DB_PATH = root / "test.sqlite3"
        chatlog._DB_PATH = docops._DB_PATH
        docops.DOCS_DIR = root / "documents"
        docops.EXPORTS_DIR = docops.DOCS_DIR / "exports"
        docops.ARTIFACTS_DIR = docops.EXPORTS_DIR / ".artifacts"
        docops.DOCS_DIR.mkdir()
        docops.EXPORTS_DIR.mkdir()
        docops.ARTIFACTS_DIR.mkdir()
        self.env = patch.dict(
            os.environ,
            {
                "PJ_TOOL_BRIDGE_TOKEN": "bridge-secret",
                "OPENAI_API_KEY": "test-key",
            },
            clear=False,
        )
        self.env.start()
        self.client = realtime_server.app.test_client()
        self.auth = {"Authorization": "Bearer bridge-secret"}

    def tearDown(self):
        self.env.stop()
        for name, value in self.old_docops.items():
            setattr(docops, name, value)
        chatlog._DB_PATH = self.old_chat_db
        self.temp_dir.cleanup()

    @staticmethod
    def _sections():
        return json.dumps({
            "Attendees": "PJ and the owner",
            "Context": "Document download validation",
            "Discussion": "Every created document must be downloadable.",
            "Decisions": "Use immutable chat artifacts.",
            "Action Items": "Ship after validation.",
        })

    def _draft(self, *, finalize=False):
        result = docops.draft_document(
            "meeting_memo",
            "Download Validation",
            self._sections(),
            finalize=finalize,
        )
        self.assertNotIn("artifact_error", result)
        self.assertEqual(result["artifact"]["status"], "ready")
        return result

    def test_created_document_has_immutable_markdown_artifact(self):
        drafted = self._draft()
        artifact = drafted["artifact"]
        self.assertEqual(artifact["format"], "md")
        self.assertFalse(artifact["audience_ready"])
        self.assertNotIn("path", artifact)

        resolved = docops.resolve_export_artifact(artifact["artifact_id"])
        self.assertEqual(resolved, artifact | {"created_at": resolved["created_at"]})
        metadata, snapshot = docops.open_export_artifact_snapshot(
            artifact["artifact_id"]
        )
        try:
            self.assertEqual(metadata["sha256"], artifact["sha256"])
            self.assertIn(b"# Download Validation", snapshot.read())
        finally:
            snapshot.close()

    def test_finalization_creates_new_audience_ready_artifact(self):
        drafted = self._draft()
        finalized = docops.finalize_document(drafted["doc_id"])
        self.assertEqual(finalized["status"], "final")
        self.assertTrue(finalized["artifact"]["audience_ready"])
        self.assertNotEqual(
            drafted["artifact"]["artifact_id"],
            finalized["artifact"]["artifact_id"],
        )
        repeated = docops.finalize_document(drafted["doc_id"])
        self.assertEqual(repeated["status"], "already_final")
        self.assertEqual(
            repeated["artifact"]["artifact_id"],
            finalized["artifact"]["artifact_id"],
        )

    def test_all_supported_exports_are_registered(self):
        drafted = self._draft(finalize=True)
        formats = ["html", "pdf", "pptx", "xlsx"]
        if shutil.which("textutil"):
            formats.extend(["docx", "rtf"])
        expected_mime_types = {
            "html": "text/html; charset=utf-8",
            "pdf": "application/pdf",
            "docx": (
                "application/vnd.openxmlformats-officedocument."
                "wordprocessingml.document"
            ),
            "rtf": "application/rtf",
            "pptx": (
                "application/vnd.openxmlformats-officedocument."
                "presentationml.presentation"
            ),
            "xlsx": (
                "application/vnd.openxmlformats-officedocument."
                "spreadsheetml.sheet"
            ),
        }
        for format_name in formats:
            with self.subTest(format=format_name):
                exported = docops.export_document(
                    drafted["doc_id"], format=format_name
                )
                self.assertNotIn("error", exported)
                artifact = exported["artifact"]
                self.assertEqual(artifact["format"], format_name)
                self.assertTrue(artifact["audience_ready"])
                self.assertGreater(artifact["byte_size"], 0)
                self.assertEqual(artifact["mime_type"],
                                 expected_mime_types[format_name])
                internal = docops.resolve_export_artifact(
                    artifact["artifact_id"], include_path=True
                )
                self.assertEqual(internal["status"], "ready")
                content = Path(internal["path"]).read_bytes()
                if format_name == "pdf":
                    self.assertTrue(content.startswith(b"%PDF-"))
                elif format_name in ("pptx", "xlsx"):
                    self.assertTrue(content.startswith(b"PK"))

                downloaded = self.client.get(
                    artifact["download_url"], headers=self.auth
                )
                try:
                    self.assertEqual(downloaded.status_code, 200)
                    self.assertEqual(downloaded.data, content)
                    self.assertEqual(
                        downloaded.headers["Content-Type"],
                        expected_mime_types[format_name],
                    )
                finally:
                    downloaded.close()

                if format_name == "pptx":
                    presentation = Presentation(internal["path"])
                    self.assertGreaterEqual(len(presentation.slides), 2)
                    self.assertEqual(
                        presentation.slides[0].shapes[4].text,
                        "Download Validation",
                    )
                elif format_name == "xlsx":
                    workbook = load_workbook(
                        internal["path"], read_only=True
                    )
                    try:
                        sheet = workbook["Document"]
                        self.assertEqual(
                            sheet["A1"].value, "Download Validation"
                        )
                        self.assertEqual(sheet["A8"].value, "Attendees")
                    finally:
                        workbook.close()

    def test_tampered_immutable_artifact_is_rejected(self):
        artifact = self._draft()["artifact"]
        internal = docops.resolve_export_artifact(
            artifact["artifact_id"], include_path=True
        )
        Path(internal["path"]).write_bytes(b"tampered")
        blocked = docops.resolve_export_artifact(artifact["artifact_id"])
        self.assertEqual(blocked["status"], "blocked")

        response = self.client.get(
            f"/responses/artifacts/{artifact['artifact_id']}",
            headers=self.auth,
        )
        self.assertEqual(response.status_code, 409)
        rendered = response.get_data(as_text=True)
        self.assertNotIn(str(self.temp_dir.name), rendered)
        self.assertNotIn(".artifacts", rendered)

    def test_resume_restores_artifact_and_authenticated_download(self):
        artifact = self._draft(finalize=True)["artifact"]
        session = chatlog.new_session(channel="web")
        self.assertTrue(
            chatlog.link_session_artifact(session["id"], artifact["artifact_id"])
        )

        resumed = self.client.post(
            f"/responses/sessions/{session['id']}/resume",
            json={},
            headers=self.auth,
        )
        self.assertEqual(resumed.status_code, 200)
        self.assertEqual(
            resumed.get_json()["session"]["artifacts"][0]["artifact_id"],
            artifact["artifact_id"],
        )

        denied = self.client.get(
            f"/responses/artifacts/{artifact['artifact_id']}"
        )
        self.assertEqual(denied.status_code, 401)
        downloaded = self.client.get(
            f"/responses/artifacts/{artifact['artifact_id']}",
            headers=self.auth,
        )
        try:
            self.assertEqual(downloaded.status_code, 200)
            self.assertEqual(downloaded.data.hex(), Path(
                docops.resolve_export_artifact(
                    artifact["artifact_id"], include_path=True
                )["path"]
            ).read_bytes().hex())
            self.assertIn(
                artifact["filename"],
                downloaded.headers["Content-Disposition"],
            )
            self.assertEqual(
                downloaded.headers["Cache-Control"], "private, no-store"
            )
        finally:
            downloaded.close()

    def test_artifact_event_is_linked_to_active_chat(self):
        artifact = self._draft()["artifact"]
        session = chatlog.new_session(channel="web")

        class FakeOrchestrator:
            def __init__(self, *_args, **_kwargs):
                pass

            def stream_turn(self, *_args, **_kwargs):
                yield {"type": "artifact.ready", **artifact}
                yield {
                    "type": "completion",
                    "text": "Document ready.",
                    "artifacts": [artifact],
                    "_response_id": "resp-document",
                }

        with patch.object(
            realtime_server, "ResponsesOrchestrator", FakeOrchestrator
        ):
            response = self.client.post(
                f"/responses/sessions/{session['id']}/turns",
                json={"message": "Create a document"},
                headers=self.auth,
            )
            self.assertEqual(response.status_code, 200)
            self.assertIn(b'"type": "artifact.ready"', response.data)

        self.assertEqual(
            chatlog.list_session_artifact_ids(session["id"]),
            [artifact["artifact_id"]],
        )

    def test_tool_and_stream_errors_redact_embedded_server_paths(self):
        with patch.object(
            realtime_server,
            "dispatch_realtime_function",
            side_effect=PermissionError(
                "cannot write /Users/private/documents/report.md"
            ),
        ):
            response = self.client.post(
                "/execute-tool",
                json={"name": "draft_document", "arguments": {}},
                headers=self.auth,
            )
        self.assertEqual(response.status_code, 500)
        self.assertNotIn("/Users/private", response.get_data(as_text=True))
        self.assertIn("[server path redacted]", response.get_data(as_text=True))

        session = chatlog.new_session(channel="web")

        class FailingOrchestrator:
            def __init__(self, *_args, **_kwargs):
                pass

            def stream_turn(self, *_args, **_kwargs):
                raise PermissionError(
                    "cannot open /Users/private/documents/report.md"
                )
                yield

        with (
            patch.object(
                realtime_server, "ResponsesOrchestrator", FailingOrchestrator
            ),
            patch.object(
                realtime_server, "OPENAI_CLIENT_FACTORY", return_value=object()
            ),
        ):
            streamed = self.client.post(
                f"/responses/sessions/{session['id']}/turns",
                json={"message": "Create a document"},
                headers=self.auth,
            )
            rendered = streamed.get_data(as_text=True)
        self.assertEqual(streamed.status_code, 200)
        self.assertNotIn("/Users/private", rendered)
        self.assertIn("[server path redacted]", rendered)


if __name__ == "__main__":
    unittest.main()
